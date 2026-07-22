#!/usr/bin/env python3
"""师助AI - PPT生成 API (部署到 Railway / Render)"""
import os, json, io, base64, http.server, uuid, tempfile, urllib.request, urllib.error, urllib.parse, textwrap, subprocess
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

DL_DIR = tempfile.mkdtemp(prefix="seat_dl_")
PPT_PIPELINE_VERSION = 'native-pptxgen-v1'

THEME = {
    'exam':    ('#1A237E', '#2B579A', '#E53935', '#E8EAF6'),
    'safety':  ('#BF360C', '#E65100', '#FF6F00', '#FFF3E0'),
    'holiday': ('#880E4F', '#AD1457', '#D81B60', '#FCE4EC'),
    'health':  ('#1B5E20', '#2E7D32', '#43A047', '#E8F5E9'),
    'plan':    ('#4A148C', '#6A1B9A', '#7B1FA2', '#F3E5F5'),
    'default': ('#283593', '#3949AB', '#4A6CF7', '#E8EAF6'),
}

def rgb(h):
    return RGBColor(int(h[1:3],16), int(h[3:5],16), int(h[5:7],16))

def detect(t):
    s = (t or '').lower()
    if '考试' in s or '期中' in s or '复习' in s: return 'exam'
    if '安全' in s or '消防' in s or '防溺水' in s: return 'safety'
    if '节日' in s or '元旦' in s or '国庆' in s: return 'holiday'
    if '健康' in s or '环保' in s or '运动' in s: return 'health'
    if '总结' in s or '计划' in s or '规划' in s: return 'plan'
    return 'default'

def generate_image_url(prompt, api_key):
    """调用阿里 Z-Image 同步文生图接口，返回临时图片 URL。"""
    image_prompt = '面向中国中小学课堂的教育教学PPT插图，人物和场景符合中国校园语境，不要正面对镜头摆拍；画面绝对不要任何文字、汉字、数字、字母、标志、海报或水印；构图简洁、主体清晰：' + prompt
    payload = json.dumps({
        'model': 'z-image-turbo',
        'input': {'messages': [{'role': 'user', 'content': [{'text': image_prompt}]}]}
    }).encode()
    req = urllib.request.Request(
        'https://ws-5ol6m5p8f4hikz1a.cn-beijing.maas.aliyuncs.com/api/v1/services/aigc/multimodal-generation/generation',
        data=payload,
        headers={'Authorization': 'Bearer ' + api_key, 'Content-Type': 'application/json'}
    )
    result = json.loads(urllib.request.urlopen(req, timeout=45).read().decode('utf-8'))
    if result.get('code'):
        raise RuntimeError('阿里生图失败：' + result.get('message', result['code']))
    content = result.get('output', {}).get('choices', [{}])[0].get('message', {}).get('content', [])
    for item in content:
        if isinstance(item, dict) and item.get('image'):
            return item['image']
    raise RuntimeError('阿里生图未返回图片地址')

def search_image_url(query, api_key, timeout=15):
    """使用阿里文搜图获取真实图片地址，适用于人物、事件、器材等事实性页面。"""
    payload = json.dumps({
        'model': 'qwen3.6-flash',
        'input': '请搜索一张可用于教学PPT的真实图片：' + query,
        'tools': [{'type': 'web_search_image'}],
        'store': False
    }).encode()
    req = urllib.request.Request(
        'https://ws-5ol6m5p8f4hikz1a.cn-beijing.maas.aliyuncs.com/compatible-mode/v1/responses',
        data=payload,
        headers={'Authorization': 'Bearer ' + api_key, 'Content-Type': 'application/json'}
    )
    response = json.loads(urllib.request.urlopen(req, timeout=timeout).read().decode('utf-8'))
    for item in response.get('output', []):
        if item.get('type') != 'web_search_image_call':
            continue
        try:
            images = json.loads(item.get('output', '[]'))
            if images and images[0].get('url'):
                return images[0]['url']
        except (TypeError, ValueError):
            continue
    raise RuntimeError('联网搜图未返回可下载图片')

def source_page_image_url(page_urls, timeout=12):
    """优先取新闻/官方网页公开的 og:image，而不是把同一主题重新画成 AI 图。"""
    import re
    from html import unescape
    for page_url in page_urls or []:
        try:
            request = urllib.request.Request(page_url, headers={'User-Agent': 'Mozilla/5.0'})
            with urllib.request.urlopen(request, timeout=timeout) as response:
                html = response.read(1500000).decode('utf-8', errors='ignore')
            for tag in re.findall(r'<meta\b[^>]*>', html, flags=re.I):
                attrs = dict((k.lower(), unescape(v)) for k, _, v in re.findall(r'''([\w:-]+)\s*=\s*(["'])(.*?)\2''', tag, flags=re.S))
                marker = (attrs.get('property') or attrs.get('name') or '').lower()
                image_url = attrs.get('content', '').strip()
                if marker in ('og:image', 'twitter:image', 'twitter:image:src') and image_url:
                    return urllib.parse.urljoin(page_url, image_url)
        except Exception:
            continue
    raise RuntimeError('来源网页未提供可用公开配图')

def build_ppt_research(topic, api_key, detail='中'):
    """只获取真实检索来源；绝不把模型自行写出的“来源”当作事实。"""
    topic = str(topic or '').strip()
    if not topic:
        raise RuntimeError('PPT主题不能为空')
    depth_hint = {
        '短': '只检索最必要的官方规范、权威说明或典型资料。',
        '中': '检索能够支撑课堂问题、解释和做法的资料。',
        '长': '在核心资料外补充可用于讨论、练习或延伸的资料。',
    }.get(str(detail), '检索能够支撑课堂问题、解释和做法的资料。')
    prompt = '''请联网为中国中小学课堂备课检索主题“%s”。%s
优先官方部门、学校教育机构、专业机构和主流媒体的原始报道；不要编造事件、数字、文件或来源。请写一份简短研究摘要，说明哪些材料适合课堂使用、哪些事实需要谨慎表述。''' % (topic, depth_hint)
    payload = json.dumps({
        'model': 'qwen-plus',
        'input': {'messages': [{'role': 'user', 'content': prompt}]},
        'parameters': {
            'enable_search': True,
            'search_options': {'forced_search': True, 'enable_source': True},
            'result_format': 'message',
            'temperature': 0.2,
        },
    }).encode()
    req = urllib.request.Request(
        'https://ws-5ol6m5p8f4hikz1a.cn-beijing.maas.aliyuncs.com/api/v1/services/aigc/text-generation/generation',
        data=payload,
        headers={'Authorization': 'Bearer ' + api_key, 'Content-Type': 'application/json'}
    )
    # 研究是生成的硬门槛，但不能在网络异常时让用户无止境等待。
    result = json.loads(urllib.request.urlopen(req, timeout=45).read().decode('utf-8'))
    output = result.get('output', {})
    raw = output.get('choices', [{}])[0].get('message', {}).get('content', '')
    # 来源必须来自阿里真实搜索结果，不能接受模型在 JSON 里自行编写的 URL。
    real_sources = []
    for source in output.get('search_info', {}).get('search_results', []):
        url = str(source.get('url', '')).strip()
        title = str(source.get('title', '')).strip()
        if url.startswith(('https://', 'http://')) and title:
            real_sources.append({
                'title': title,
                'url': url,
                'excerpt': str(source.get('snippet', source.get('content', ''))).strip()[:500],
            })
    def source_score(source):
        domain = urllib.parse.urlparse(source['url']).netloc.lower()
        title = source['title']
        if domain.endswith('.gov.cn') or domain.endswith('.edu.cn'):
            return 100
        if any(name in domain for name in ('news.cn', 'xinhuanet.com', 'people.com.cn', 'cctv.com', 'chinanews.com')):
            return 75
        if any(name in title for name in ('教育部', '应急管理部', '国家', '中国教育报', '新华社', '人民日报')):
            return 65
        return 10
    real_sources.sort(key=source_score, reverse=True)
    if len(real_sources) < 2:
        raise RuntimeError('联网研究未获得至少2个可核验来源')
    if source_score(real_sources[0]) < 65:
        raise RuntimeError('联网研究未找到可信官方或权威来源，已停止生成')
    return {
        'research_note': str(raw).strip()[:3000],
        'sources': real_sources[:4],
        'assets': [],
    }

def _json_from_model(raw):
    """兼容模型偶尔用 Markdown 代码块包裹 JSON 的情况。"""
    import re
    match = re.search(r'\{[\s\S]*\}', str(raw or ''))
    if not match:
        raise RuntimeError('故事板未返回 JSON')
    try:
        return json.loads(match.group(0))
    except ValueError as err:
        raise RuntimeError('故事板 JSON 格式错误') from err

def _qwen_reply(api_key, messages, temperature=0.4, timeout=45):
    payload = json.dumps({
        'model': 'qwen-plus', 'messages': messages, 'temperature': temperature,
        'max_tokens': 4096, 'stream': False,
    }).encode()
    req = urllib.request.Request(
        'https://ws-5ol6m5p8f4hikz1a.cn-beijing.maas.aliyuncs.com/compatible-mode/v1/chat/completions',
        data=payload,
        headers={'Authorization': 'Bearer ' + api_key, 'Content-Type': 'application/json'}
    )
    result = json.loads(urllib.request.urlopen(req, timeout=timeout).read().decode('utf-8'))
    reply = result.get('choices', [{}])[0].get('message', {}).get('content', '')
    if not str(reply).strip():
        raise RuntimeError('阿里模型未返回内容')
    return str(reply).strip()

def _outline_to_slides(outline):
    """把自由 Markdown 大纲转成页面；模型只负责内容，程序才负责结构。"""
    import re
    slides, current = [], None
    subtitle = ''
    for raw_line in str(outline or '').splitlines():
        line = raw_line.strip()
        if not line:
            continue
        if line.startswith('# ') and not subtitle:
            subtitle = line[2:].strip()
            continue
        if re.match(r'^#{2,3}\s+', line):
            if current and current['content']:
                slides.append(current)
            current = {'title': re.sub(r'^#{2,3}\s+', '', line).strip(), 'content': []}
            continue
        if current:
            line = re.sub(r'^[-*•▪]\s*', '', line)
            line = re.sub(r'^\d+[、.．]\s*', '', line)
            if line:
                current['content'].append(line)
    if current and current['content']:
        slides.append(current)
    return subtitle, slides

def _qwen_search_storyboard(api_key, prompt, timeout=75):
    """一次阿里联网请求：返回逐页 JSON 和该请求真实返回的来源。"""
    payload = json.dumps({
        'model': 'qwen-plus',
        'input': {'messages': [{'role': 'user', 'content': prompt}]},
        'parameters': {
            'enable_search': True,
            'search_options': {'forced_search': True, 'enable_source': True},
            'result_format': 'message',
            'temperature': 0.45,
        },
    }).encode()
    req = urllib.request.Request(
        'https://ws-5ol6m5p8f4hikz1a.cn-beijing.maas.aliyuncs.com/api/v1/services/aigc/text-generation/generation',
        data=payload,
        headers={'Authorization': 'Bearer ' + api_key, 'Content-Type': 'application/json'}
    )
    result = json.loads(urllib.request.urlopen(req, timeout=timeout).read().decode('utf-8'))
    output = result.get('output', {})
    raw = str(output.get('choices', [{}])[0].get('message', {}).get('content', '')).strip()
    if not raw:
        raise RuntimeError('阿里联网模型未返回故事板')
    sources = []
    for source in output.get('search_info', {}).get('search_results', []):
        url = str(source.get('url', '')).strip()
        title = str(source.get('title', '')).strip()
        if url.startswith(('https://', 'http://')) and title:
            sources.append({
                'title': title,
                'url': url,
                'excerpt': str(source.get('snippet', source.get('content', ''))).strip()[:500],
            })
    # 去重但保持搜索返回的顺序，资料页只展示实际搜索结果。
    unique = []
    seen = set()
    for source in sources:
        if source['url'] not in seen:
            seen.add(source['url'])
            unique.append(source)
    if len(unique) < 2:
        raise RuntimeError('联网故事板未获得至少2个可核验来源')
    return raw, unique[:6]

def build_ppt_deck(body, api_key):
    """PPT 的完整链路：一次联网研究并直接形成逐页故事板 → 后端排版。"""
    topic = str(body.get('topic') or body.get('title') or '').strip()
    detail = str(body.get('detail') or body.get('wordCount') or '中')
    debug = {}
    def trace(stage, value):
        if body.get('debug'):
            debug[stage] = value
            print('[PPT TRACE] ' + stage + ': ' + json.dumps(value, ensure_ascii=False)[:12000])
    length_hint = {
        '短': '4-6 页内容页', '中': '6-9 页内容页', '长': '9-13 页内容页',
    }.get(detail, '6-9 页内容页')
    grade = str(body.get('grade') or '').strip()
    subject = str(body.get('subject') or '').strip()
    academic_words = ('考试', '期中', '期末', '复习', '学习', '学法', '动员', '成绩', '备考', '中考', '高考', '学科', '阅读', '数学', '语文', '英语')
    topic_is_academic = any(word in topic for word in academic_words)
    user_context = '适用对象：%s。' % ((grade + '学生') if grade else '中小学生')
    if topic_is_academic and subject:
        user_context += '这是与学习有关的班会，可自然结合%s学科的学习情境，但不得生硬类比。' % subject
    extra = str(body.get('userDetail') or '').strip()[:1200]
    storyboard_prompt = '''请联网检索，并直接为中国中小学班会“%s”制作逐页 PPT 故事板。不要先写讲稿、再压缩成大纲；请在这一次回答里同时完成资料理解、教学设计与逐页内容。

%s
教师补充要求：%s
内容页数量：%s。封面、资料来源、结尾由程序添加，你不要生成它们。

必须只输出合法 JSON，不要 Markdown 代码块：
{"slides":[{"title":"页面标题","type":"narrative|explain|steps|scenario|case","content":["本页正文或要点"],"image_query":"真实图片检索词或空字符串","source_ids":["S1"]}]}

写作准则：
1. 以你自己的教学设计能力为主：课堂逻辑、解释、启发式提问、假如情境、活动设计和行动建议都应由你主动组织，不要把检索资料拼贴成课件。联网只用于校验或补充少数必须真实的内容：政策要求、精确数据、近期/具体事件、可引用资源，以及真实配图。不要把资料标题当成资料正文；没有可靠来源支撑时，不写精确数字、秒数、年份、具体事故、人物、地点、引语或学校安排。
2. 每页只讲一个完整问题，顺序应像一堂真实班会自然推进：情境或问题→解释/辨析→具体做法→讨论或行动；不要机械套模板，也不要重复“防溺水很重要”。
3. narrative/case 页可放 120-260 字的完整段落；explain 页放 2-4 条有解释力的内容；steps 页放 3-6 个可执行的短步骤；scenario 页只放一个可讨论的情境。禁止把一句话拆成多页或使用“续页”。
4. 不得冒充该校教师，不得虚构本校经历、学生原话、新闻事件；课堂情境必须用“假如/例如”明确标注。
5. 对叙事、解释、案例等适合视觉呈现的页面，尽量填写 image_query，写成可联网搜索的具体中文真实照片词（例如“中小学生 防溺水 安全教育 宣传活动 真实照片”）；步骤、讨论等不需要配图的页面可留空。不要要求 AI 图，不要让图片含文字。
6. source_ids 只能填写本次联网检索实际返回来源的 S 编号；有来源的页面会优先尝试使用该来源网页公开配图，没有对应来源则为空数组。''' % (topic, user_context, extra or '无', length_hint)
    trace('01_storyboard_request', {'detail': detail, 'length_hint': length_hint, 'topic': topic})
    try:
        raw, sources = _qwen_search_storyboard(api_key, storyboard_prompt, timeout=75)
    except Exception as e:
        trace('01_storyboard_error', {'type': type(e).__name__, 'message': str(e)[:1000]})
        raise RuntimeError('PPT 联网故事板生成失败：' + str(e)[:500]) from e
    trace('02_research_sources', sources)
    trace('03_storyboard_json', raw)
    plan = _json_from_model(raw)
    outline_slides = plan.get('slides', []) if isinstance(plan, dict) else []
    trace('04_storyboard_parsed', outline_slides)
    slides = []
    allowed_kinds = {'narrative', 'explain', 'steps', 'scenario', 'case', 'compare'}
    for item in outline_slides if isinstance(outline_slides, list) else []:
        if not isinstance(item, dict):
            continue
        title = str(item.get('title') or '').strip()[:30]
        content = [str(x).strip() for x in item.get('content', []) if str(x).strip()]
        if not title or not content:
            continue
        kind = str(item.get('type') or 'explain').lower()
        kind = kind if kind in allowed_kinds else 'explain'
        source_titles = []
        for source_id in item.get('source_ids', []) if isinstance(item.get('source_ids'), list) else []:
            if isinstance(source_id, str) and source_id.startswith('S') and source_id[1:].isdigit():
                index = int(source_id[1:]) - 1
                if 0 <= index < len(sources):
                    source_titles.append(sources[index]['title'])
        slides.append({'type': kind, 'title': title, 'content': content, 'source_titles': source_titles[:2], 'source_urls': [sources[int(sid[1:]) - 1]['url'] for sid in item.get('source_ids', []) if isinstance(sid, str) and sid.startswith('S') and sid[1:].isdigit() and 0 <= int(sid[1:]) - 1 < len(sources)][:2], 'image_query': str(item.get('image_query') or '').strip()[:140]})
    if len(slides) < 2:
        raise RuntimeError('故事板未形成足够的可核验内容，已停止生成')
    max_middle = {'短': 5, '中': 9, '长': 13}.get(detail, 9)
    slides = slides[:max_middle]
    trace('05_render_slides_before_images', slides)
    # 图片只能由故事板主动指定；不再按标题盲猜搜索词，避免出现无关摆拍图。
    assets = []
    image_limit = {'短': 3, '中': 5, '长': 7}.get(detail, 5)
    visual_candidates = [
        (index, slide) for index, slide in enumerate(slides)
        if slide.get('type') in ('narrative', 'explain', 'case') and (slide.get('image_query') or slide.get('source_urls'))
    ][:image_limit] if not body.get('_skip_assets') else []
    if visual_candidates:
        # 先取资料原网页公开配图；没有时才搜图。依次请求避免阿里搜图并发时一起超时。
        for index, slide in visual_candidates:
            try:
                source_urls = slide.get('source_urls', [])
                try:
                    trace('06_image_request', {'slide': index + 1, 'mode': 'source_page', 'sources': source_urls})
                    url = source_page_image_url(source_urls, timeout=12)
                except Exception:
                    if not slide.get('image_query'):
                        raise
                    trace('06_image_request', {'slide': index + 1, 'mode': 'image_search', 'query': slide['image_query']})
                    url = search_image_url(slide['image_query'], api_key, timeout=20)
                if any(asset.get('url') == url for asset in assets):
                    continue
                asset_id = 'A' + str(len(assets) + 1)
                assets.append({'id': asset_id, 'url': url, 'caption': slides[index]['image_query']})
                slides[index]['visual'] = {'mode': 'asset', 'asset_id': asset_id}
            except Exception as err:
                print('[PPT API] research visual skipped:', str(err)[:120])
                if body.get('debug'):
                    debug.setdefault('06_image_errors', []).append(str(err)[:300])
    trace('06_image_assets', assets)
    deck = {
        'title': topic[:40],
        'theme': body.get('theme') or detect(topic),
        'slides': ([{'type': 'cover', 'title': topic, 'content': ['基于联网资料与课堂常识整理']}]
                   + slides
                   + [{'type': 'sources', 'title': '资料来源', 'content': [(urllib.parse.urlparse(s['url']).netloc + '｜' + s['title'])[:80] for s in sources[:4]], 'source_urls': [s['url'] for s in sources[:4]]},
                      {'type': 'closing', 'title': '把安全与行动带回日常', 'content': ['课堂讨论后，请把今天学到的做法落实到具体场景。']}]),
        'assets': assets,
        'debug': debug,
    }
    return deck

def _official_ppt_config():
    """阿里官方 PPT 服务使用云账号 AK，不与百炼模型 API Key 混用。"""
    access_key_id = (os.environ.get('ALIYUN_ACCESS_KEY_ID') or os.environ.get('ALIBABA_CLOUD_ACCESS_KEY_ID') or '').strip()
    access_key_secret = (os.environ.get('ALIYUN_ACCESS_KEY_SECRET') or os.environ.get('ALIBABA_CLOUD_ACCESS_KEY_SECRET') or '').strip()
    workspace_id = (os.environ.get('BAILIAN_WORKSPACE_ID') or '').strip()
    missing = []
    if not access_key_id: missing.append('ALIYUN_ACCESS_KEY_ID')
    if not access_key_secret: missing.append('ALIYUN_ACCESS_KEY_SECRET')
    if not workspace_id: missing.append('BAILIAN_WORKSPACE_ID')
    if missing:
        raise RuntimeError('官方PPT服务尚未配置：请在 Railway 添加 ' + '、'.join(missing))
    return access_key_id, access_key_secret, workspace_id

def _deck_to_official_outline(deck):
    """把已审过的课堂故事板转成官方 PPT 服务要求的 Markdown 大纲。"""
    lines = ['# ' + str(deck.get('title') or '教学课件')]
    for slide in deck.get('slides', []):
        kind = str(slide.get('type') or '')
        if kind in ('cover', 'sources', 'closing'):
            continue
        title = str(slide.get('title') or '').strip()
        content = [str(x).strip() for x in slide.get('content', []) if str(x).strip()]
        if not title or not content:
            continue
        lines.append('## ' + title)
        for item in content:
            lines.append('- ' + item)
    if len(lines) < 3:
        raise RuntimeError('PPT大纲内容不足，未提交官方生成服务')
    return '\n'.join(lines)

def _official_get_ppt_info(client, workspace_id, task_id):
    """官方 Demo 中的 GetPptInfo 轮询接口。"""
    from alibabacloud_aimiaobi20230801 import models as aimiaobi_models
    from darabonba.runtime import RuntimeOptions
    request = aimiaobi_models.GetPptInfoRequest(workspace_id=workspace_id, task_id=task_id)
    response = client.get_ppt_info_with_options(
        request, RuntimeOptions(read_timeout=30000, connect_timeout=10000)
    )
    return getattr(getattr(response, 'body', None), 'data', None)

def build_official_ppt(body, api_key):
    """内容大纲由模型产生，版式、素材和 PPTX 由阿里官方 PPT 服务生成。"""
    access_key_id, access_key_secret, workspace_id = _official_ppt_config()
    try:
        from alibabacloud_tea_openapi import models as open_api_models
        from alibabacloud_aimiaobi20230801.client import Client as AimiaobiClient
        from alibabacloud_aimiaobi20230801 import models as aimiaobi_models
        from darabonba.runtime import RuntimeOptions
    except ImportError as err:
        raise RuntimeError('官方PPT依赖未安装，请确认 Railway 已按最新 requirements.txt 重新部署') from err

    outline_body = dict(body)
    outline_body['_skip_assets'] = True
    deck = build_ppt_deck(outline_body, api_key)
    outline = _deck_to_official_outline(deck)
    if body.get('debug'):
        print('[PPT TRACE] 07_official_outline: ' + outline[:12000])

    config = open_api_models.Config(
        access_key_id=access_key_id,
        access_key_secret=access_key_secret,
        region_id='cn-beijing',
        endpoint='aimiaobi.cn-beijing.aliyuncs.com',
    )
    client = AimiaobiClient(config)
    task_id = str(uuid.uuid4())
    request = aimiaobi_models.InitiatePptCreationV2Request(
        workspace_id=workspace_id,
        task_id=task_id,
        outline=outline,
        process_type=4,
    )
    response = client.initiate_ppt_creation_v2with_options(
        request, RuntimeOptions(read_timeout=30000, connect_timeout=10000)
    )
    data = getattr(getattr(response, 'body', None), 'data', None)
    if not data:
        raise RuntimeError('官方PPT服务未返回创建任务')
    if body.get('debug'):
        print('[PPT TRACE] 08_official_task: ' + json.dumps({
            'task_id': task_id,
            'export_task_id': getattr(data, 'export_task_id', None),
        }, ensure_ascii=False))

    # 官方场景五会在任务信息中返回最终链接；最长约两分钟，避免小程序无限等待。
    import time
    for attempt in range(40):
        info = _official_get_ppt_info(client, workspace_id, task_id)
        links = getattr(info, 'export_file_link', None) or []
        if isinstance(links, str):
            links = [links]
        if links:
            download_url = str(links[0])
            image_request = urllib.request.Request(download_url, headers={'User-Agent': 'Mozilla/5.0'})
            pptx_bytes = urllib.request.urlopen(image_request, timeout=45).read()
            if len(pptx_bytes) < 5000:
                raise RuntimeError('官方PPT导出文件异常')
            if body.get('debug'):
                print('[PPT TRACE] 09_official_export: ' + json.dumps({'attempt': attempt + 1, 'bytes': len(pptx_bytes)}, ensure_ascii=False))
            return pptx_bytes, deck
        time.sleep(3)
    raise RuntimeError('官方PPT仍在生成，请稍后重试')

def fit_slides(raw_slides):
    """把过长要点拆成续页，确保每页最多 4 条、每条最多约 42 个字符。"""
    if len(raw_slides) <= 2:
        return raw_slides
    fitted = [raw_slides[0]]
    for source in raw_slides[1:-1]:
        kind = str(source.get('type', '')).lower()
        items = []
        for item in source.get('content', []):
            item = str(item).strip()
            if not item:
                continue
            # 案例/叙事页允许完整段落，不把一个事件或讲述机械切成多条短句。
            while kind not in ('case', 'narrative') and len(item) > 42:
                cut = max(item.rfind('，', 0, 42), item.rfind('。', 0, 42), item.rfind('；', 0, 42), item.rfind('、', 0, 42))
                cut = cut + 1 if cut >= 18 else 42
                items.append(item[:cut])
                item = item[cut:].lstrip('，。；、 ')
            if item:
                items.append(item)
        # 互动和步骤页需要留白；事实、案例、做法页可以承载更多信息。
        max_items = {
            'scenario': 3, 'steps': 6, 'compare': 6,
            'fact': 6, 'case': 6, 'narrative': 1, 'action': 6, 'sources': 4,
        }.get(kind, 5)
        chunks = [items[i:i + max_items] for i in range(0, len(items), max_items)] or [[]]
        for chunk_index, chunk in enumerate(chunks):
            slide = dict(source)
            slide['content'] = chunk
            if chunk_index:
                slide['title'] = str(source.get('title', '')) + '（续）'
            fitted.append(slide)
    fitted.append(raw_slides[-1])
    return fitted

def gen(data):
    th = THEME.get(data.get('theme') or detect(data.get('title','')), THEME['default'])
    D = rgb(th[0]); P = rgb(th[1]); A = rgb(th[2]); L = rgb(th[3])
    W = RGBColor(0xFF,0xFF,0xFF); T = RGBColor(0x33,0x33,0x33)
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    slides = [dict(s, content=list(s.get('content', []))) for s in data.get('slides', [])]
    asset_map = {
        str(asset.get('id')): asset for asset in data.get('assets', [])
        if isinstance(asset, dict) and asset.get('id') and asset.get('url')
    }
    # 先移除插图标记，避免它作为文字渲染到 PPT 上。每份 PPT 最多生成两张图，
    # 控制等待时间和费用；提示词也会要求模型最多标记两页。
    import re
    image_requests = []
    for source in slides:
        clean_content = []
        seen_items = set()
        visual = source.get('visual') if isinstance(source.get('visual'), dict) else {}
        if visual.get('mode') == 'asset' and str(visual.get('asset_id', '')) in asset_map:
            asset = asset_map[str(visual['asset_id'])]
            source['_image_mode'] = '素材'
            source['_image_url'] = asset['url']
            source['_image_prompt'] = asset.get('caption', '')
            image_requests.append(source)
        elif visual.get('mode') in ('search', 'generate') and str(visual.get('prompt', '')).strip():
            source['_image_mode'] = '搜图' if visual['mode'] == 'search' else '生图'
            source['_image_prompt'] = str(visual['prompt']).strip()
            image_requests.append(source)
        for item in source.get('content', []):
            match = re.search(r'\[(\u641c\u56fe|\u751f\u56fe|\u63d2\u56fe)\uff1a(.+?)\]', str(item))
            if match and not source.get('_image_prompt'):
                source['_image_mode'] = match.group(1)
                source['_image_prompt'] = match.group(2).strip()
                item = str(item).replace(match.group(0), '').strip()
                image_requests.append(source)
            # 模型常自行加编号，版式也会加编号；统一去除，避免“1. 1.”。
            item = re.sub(r'^\s*(?:[0-9]+[、.．]|[一二三四五六七八九十]+、|[▪•\-])\s*', '', str(item)).strip()
            normalized = re.sub(r'\s+', '', item)
            if item and normalized not in seen_items:
                clean_content.append(item)
                seen_items.add(normalized)
        source['content'] = clean_content
    # 不为“看起来有图”而强行补图。只有研究素材包或故事板明确选择的画面才会进入 PPT，
    # 这样避免与主题无关的摆拍图、以及带乱码文字的生图。
    if len(image_requests) > 3:
        # 故事板可能给多页标注配图；不应因此放弃整份PPT。
        # 优先保留真实搜图和案例/事实/情境页，其他页面自然回退为纯文字版。
        priority = {'case': 0, 'fact': 1, 'scenario': 2, 'steps': 3, 'compare': 4}
        ranked = sorted(
            enumerate(image_requests),
            key=lambda pair: (
                0 if pair[1].get('_image_mode') == '搜图' else 1,
                priority.get(str(pair[1].get('type', '')).lower(), 5),
                pair[0],
            )
        )
        selected = {id(source) for _, source in ranked[:3]}
        for source in image_requests:
            if id(source) not in selected:
                source.pop('_image_mode', None)
                source.pop('_image_prompt', None)
        image_requests = [source for _, source in ranked[:3]]
    image_key = os.environ.get('AI_API_KEY', '').strip()
    for source in image_requests:
        if not image_key:
            raise RuntimeError('PPT 需要插图，但 Railway 未配置 AI_API_KEY')
        image_path = os.path.join(tempfile.mkdtemp(prefix='ppt_img_'), 'image.png')
        try:
            if source.get('_image_mode') == '素材':
                image_url = source['_image_url']
            elif source.get('_image_mode') == '搜图':
                image_url = search_image_url(source['_image_prompt'], image_key)
            else:
                image_url = generate_image_url(source['_image_prompt'], image_key)
            # 不少新闻/图片站会拒绝无 User-Agent 的服务器下载。
            image_request = urllib.request.Request(image_url, headers={'User-Agent': 'Mozilla/5.0'})
            with urllib.request.urlopen(image_request, timeout=15) as image_response, open(image_path, 'wb') as image_file:
                image_file.write(image_response.read())
            source['_image_path'] = image_path
        except Exception as image_error:
            # 外站图片可能防盗链或暂时不可访问。跳过图片即可，不能为此让正文长时间等待，
            # 更不能擅自改成可能含乱码文字的 AI 生图。
            if source.get('_image_mode') in ('搜图', '素材'):
                print('[PPT API] research image skipped:', str(image_error)[:120])
            else:
                print('[PPT API] AI image skipped:', str(image_error)[:120])
            source.pop('_image_path', None)
    slides = fit_slides(slides); n = len(slides)

    for idx, s in enumerate(slides):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # 除封面外统一页码，避免内容页出现重复或错乱编号。
        if idx > 0:
            footer = slide.shapes.add_textbox(Inches(11.8), Inches(7.03), Inches(1.0), Inches(0.25))
            fp = footer.text_frame.paragraphs[0]
            fp.text = f'{idx + 1} / {len(slides)}'
            fp.font.size = Pt(9); fp.font.color.rgb = P; fp.alignment = PP_ALIGN.RIGHT
        if idx == 0:
            slide.background.fill.solid(); slide.background.fill.fore_color.rgb = D
            tb = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10), Inches(2))
            tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.text = data.get('title',''); p.font.size = Pt(52)
            p.font.bold = True; p.font.color.rgb = W; p.alignment = PP_ALIGN.CENTER
            if s.get('content'):
                p2 = tf.add_paragraph(); p2.text = s['content'][0] or ''
                p2.font.size = Pt(22); p2.font.color.rgb = A; p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(12)
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(1.9), Inches(4), Pt(3))
            line.fill.solid(); line.fill.fore_color.rgb = A; line.line.fill.background()
            continue
        if idx == n - 1:
            slide.background.fill.solid(); slide.background.fill.fore_color.rgb = L
            tb = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(9), Inches(2))
            tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.text = s.get('title',''); p.font.size = Pt(40)
            p.font.bold = True; p.font.color.rgb = D; p.alignment = PP_ALIGN.CENTER
            if s.get('content'):
                p2 = tf.add_paragraph(); p2.text = s['content'][0]; p2.font.size = Pt(22)
                p2.font.color.rgb = A; p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(16)
            continue
        kind = str(s.get('type', '')).lower()
        items = s.get('content', [])
        # 故事板页面按教学目的渲染，避免所有内容页都是同一种项目符号列表。
        if kind == 'narrative':
            title = slide.shapes.add_textbox(Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.8))
            tp = title.text_frame.paragraphs[0]; tp.text = s.get('title', '')
            tp.font.size = Pt(30); tp.font.bold = True; tp.font.color.rgb = D
            has_image = bool(s.get('_image_path'))
            panel_x = Inches(5.05) if has_image else Inches(1.05)
            panel_w = Inches(7.1) if has_image else Inches(11.1)
            panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, panel_x, Inches(1.65), panel_w, Inches(3.85))
            panel.fill.solid(); panel.fill.fore_color.rgb = L; panel.line.fill.background()
            tf = panel.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.text = ''.join(items[:1]); p.font.size = Pt(18 if has_image else 20); p.font.color.rgb = T
            p.space_after = Pt(10)
            if has_image:
                slide.shapes.add_picture(s['_image_path'], Inches(0.95), Inches(1.75), Inches(3.65), Inches(3.65))
            source_names = s.get('source_titles') or []
            if source_names:
                foot = slide.shapes.add_textbox(Inches(1.15), Inches(5.75), Inches(10.8), Inches(0.35))
                fp = foot.text_frame.paragraphs[0]; fp.text = '资料：' + '；'.join(source_names[:2])
                fp.font.size = Pt(10); fp.font.color.rgb = P
            continue
        if kind == 'explain':
            title = slide.shapes.add_textbox(Inches(0.85), Inches(0.5), Inches(11.5), Inches(0.8))
            tp = title.text_frame.paragraphs[0]; tp.text = s.get('title', '')
            tp.font.size = Pt(30); tp.font.bold = True; tp.font.color.rgb = D
            has_image = bool(s.get('_image_path'))
            panel_x = Inches(0.95)
            panel_w = Inches(6.8) if has_image else Inches(11.2)
            panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, panel_x, Inches(1.55), panel_w, Inches(4.35))
            panel.fill.solid(); panel.fill.fore_color.rgb = L; panel.line.fill.background()
            tf = panel.text_frame; tf.word_wrap = True
            for i, item in enumerate(items[:5]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = item; p.font.size = Pt(18); p.font.color.rgb = T; p.space_after = Pt(12)
            if has_image:
                slide.shapes.add_picture(s['_image_path'], Inches(8.15), Inches(1.75), Inches(3.5), Inches(3.5))
            source_names = s.get('source_titles') or []
            if source_names:
                foot = slide.shapes.add_textbox(Inches(1.05), Inches(6.05), Inches(10.8), Inches(0.3))
                fp = foot.text_frame.paragraphs[0]; fp.text = '资料：' + '；'.join(source_names[:2])
                fp.font.size = Pt(10); fp.font.color.rgb = P
            continue
        if kind == 'scenario':
            slide.background.fill.solid(); slide.background.fill.fore_color.rgb = L
            q = slide.shapes.add_textbox(Inches(1.15), Inches(1.1), Inches(11.0), Inches(1.3))
            qp = q.text_frame.paragraphs[0]; qp.text = s.get('title', '')
            qp.font.size = Pt(34); qp.font.bold = True; qp.font.color.rgb = D; qp.alignment = PP_ALIGN.CENTER
            panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(2.6), Inches(10.3), Inches(2.8))
            panel.fill.solid(); panel.fill.fore_color.rgb = W; panel.line.color.rgb = A
            tf = panel.text_frame; tf.word_wrap = True
            for i, item in enumerate(items[:3]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = item; p.font.size = Pt(21); p.font.color.rgb = T; p.alignment = PP_ALIGN.CENTER
                p.space_after = Pt(12)
            continue
        if kind == 'case':
            title = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.8), Inches(0.8))
            tp = title.text_frame.paragraphs[0]; tp.text = s.get('title', ''); tp.font.size = Pt(30); tp.font.bold = True; tp.font.color.rgb = D
            has_image = bool(s.get('_image_path'))
            text_x = Inches(5.25) if has_image else Inches(1.1)
            text_w = Inches(6.9) if has_image else Inches(11.1)
            panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, text_x, Inches(1.55), text_w, Inches(4.5))
            panel.fill.solid(); panel.fill.fore_color.rgb = L; panel.line.fill.background()
            tf = panel.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.text = '\n'.join(items[:2]); p.font.size = Pt(18); p.font.color.rgb = T
            if has_image:
                slide.shapes.add_picture(s['_image_path'], Inches(0.85), Inches(1.7), Inches(3.85), Inches(3.85))
            continue
        if kind == 'steps':
            title = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.8), Inches(0.8))
            tp = title.text_frame.paragraphs[0]; tp.text = s.get('title', ''); tp.font.size = Pt(30); tp.font.bold = True; tp.font.color.rgb = D
            # 3×2 六宫格：例如“六不准”必须完整留在一页，而不是拆成续页。
            count = min(max(len(items), 1), 6)
            columns = 3 if count > 4 else count
            card_w = 11.4 / columns
            for i, item in enumerate(items[:6]):
                col, row = i % columns, i // columns
                x = 0.9 + col * card_w
                y = 1.55 + row * 2.35
                num = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.22), Inches(y), Inches(0.56), Inches(0.56))
                num.fill.solid(); num.fill.fore_color.rgb = A; num.line.fill.background()
                np = num.text_frame.paragraphs[0]; np.text = str(i + 1); np.font.size = Pt(16); np.font.bold = True; np.font.color.rgb = W; np.alignment = PP_ALIGN.CENTER
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y + 0.72), Inches(card_w - 0.32), Inches(1.25))
                card.fill.solid(); card.fill.fore_color.rgb = L; card.line.fill.background()
                cp = card.text_frame.paragraphs[0]; cp.text = item; cp.font.size = Pt(16); cp.font.color.rgb = D; cp.font.bold = True; cp.alignment = PP_ALIGN.CENTER
            continue
        if kind in ('compare', 'sources'):
            title = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.8), Inches(0.8))
            tp = title.text_frame.paragraphs[0]; tp.text = s.get('title', ''); tp.font.size = Pt(30); tp.font.bold = True; tp.font.color.rgb = D
            if kind == 'sources':
                for i, item in enumerate(items[:4]):
                    tb = slide.shapes.add_textbox(Inches(1.4), Inches(1.8 + i * 0.9), Inches(10.3), Inches(0.6))
                    p = tb.text_frame.paragraphs[0]; p.text = '来源  ' + item; p.font.size = Pt(18); p.font.color.rgb = T
            else:
                mid = (len(items) + 1) // 2
                for column, group in enumerate((items[:mid], items[mid:])):
                    x = 0.8 + column * 6.0
                    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.55), Inches(5.6), Inches(4.4))
                    panel.fill.solid(); panel.fill.fore_color.rgb = L; panel.line.fill.background()
                    tf = panel.text_frame; tf.word_wrap = True
                    for i, item in enumerate(group):
                        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                        p.text = item; p.font.size = Pt(17); p.font.color.rgb = D; p.space_after = Pt(14)
            continue
        lt = 1 if s.get('_image_path') else (idx - 1) % 3 + 1
        if lt == 0:
            slide.background.fill.solid(); slide.background.fill.fore_color.rgb = D
            tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
            tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.text = s.get('title',''); p.font.size = Pt(40)
            p.font.bold = True; p.font.color.rgb = W; p.alignment = PP_ALIGN.CENTER
        elif lt == 1:
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.4))
            bar.fill.solid(); bar.fill.fore_color.rgb = D; bar.line.fill.background()
            tf = bar.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.text = '  ' + s.get('title',''); p.font.size = Pt(32); p.font.bold = True; p.font.color.rgb = W
            items = s.get('content', [])
            has_image = bool(s.get('_image_path'))
            text_width = Inches(6.4) if has_image else Inches(11)
            font_size = Pt(16 if len(items) >= 4 or any(len(str(x)) > 28 for x in items) else 18)
            text_x = Inches(6.2) if has_image and idx % 2 == 0 else Inches(1.2)
            for i, item in enumerate(items):
                tb2 = slide.shapes.add_textbox(text_x, Inches(1.9+i*1.05), text_width, Inches(0.9))
                tf2 = tb2.text_frame; tf2.word_wrap = True
                p2 = tf2.paragraphs[0]; p2.text = '▪  ' + item; p2.font.size = font_size; p2.font.color.rgb = T
            if has_image:
                image_x = Inches(1.0) if idx % 2 == 0 else Inches(8.2)
                slide.shapes.add_picture(s['_image_path'], image_x, Inches(2.0), Inches(4.1), Inches(4.1))
        elif lt == 2:
            items = s.get('content',[]); mid = (len(items)+1)//2
            tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.9))
            p2 = tb2.text_frame.paragraphs[0]; p2.text = s.get('title',''); p2.font.size = Pt(30); p2.font.bold = True; p2.font.color.rgb = P
            for i, item in enumerate(items[:mid]):
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.6+i*1.0), Inches(6), Inches(0.8))
                card.fill.solid(); card.fill.fore_color.rgb = L; card.line.fill.background()
                tf2 = card.text_frame; tf2.word_wrap = True; p3 = tf2.paragraphs[0]; p3.text = item
                p3.font.size = Pt(14); p3.font.color.rgb = D; p3.font.bold = True; p3.alignment = PP_ALIGN.CENTER
            for i, item in enumerate(items[mid:]):
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6+i*1.0), Inches(6), Inches(0.8))
                card.fill.solid(); card.fill.fore_color.rgb = L; card.line.fill.background()
                tf2 = card.text_frame; tf2.word_wrap = True; p3 = tf2.paragraphs[0]; p3.text = item
                p3.font.size = Pt(14); p3.font.color.rgb = D; p3.font.bold = True; p3.alignment = PP_ALIGN.CENTER
        else:
            tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11), Inches(0.9))
            p2 = tb2.text_frame.paragraphs[0]; p2.text = s.get('title',''); p2.font.size = Pt(30); p2.font.bold = True; p2.font.color.rgb = P
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(11.5), Pt(4))
            line.fill.solid(); line.fill.fore_color.rgb = A; line.line.fill.background()
            items = s.get('content', [])
            item_step = 1.15 if len(items) >= 4 else 1.4
            item_font = Pt(16 if len(items) >= 4 else 18)
            for i, item in enumerate(items):
                y = 1.8 + i * item_step
                circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(y), Inches(0.7), Inches(0.7))
                circle.fill.solid(); circle.fill.fore_color.rgb = P; circle.line.fill.background()
                tf2 = circle.text_frame; tf2.word_wrap = False; p3 = tf2.paragraphs[0]; p3.text = str(i+1)
                p3.font.size = Pt(20); p3.font.bold = True; p3.font.color.rgb = W; p3.alignment = PP_ALIGN.CENTER
                tb3 = slide.shapes.add_textbox(Inches(2.0), Inches(y+0.05), Inches(10), Inches(0.6))
                p4 = tb3.text_frame.paragraphs[0]; p4.text = item; p4.font.size = item_font; p4.font.color.rgb = T
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()

###############################################################################
# PPT generation v1 — MiMo research agent + deterministic presentation engine
#
# The previous Qwen storyboard / official-PPT fallback / local renderer chain is
# intentionally no longer called.  It mixed three incompatible decisions and
# silently produced a local deck whenever the paid service was unavailable.
# The functions below keep one owner for content (MiMo) and one owner for
# rendering (this service), with a visible research and revision loop.
###############################################################################

def _extract_json_object(raw, label='模型'):
    """Accept JSON returned with an occasional Markdown fence, but nothing else."""
    import re
    raw = str(raw or '').strip()
    raw = re.sub(r'^```(?:json)?\s*|\s*```$', '', raw, flags=re.I)
    first = raw.find('{')
    last = raw.rfind('}')
    if first < 0 or last <= first:
        raise RuntimeError(label + '没有返回可读取的结构化方案')
    try:
        return json.loads(raw[first:last + 1])
    except ValueError as exc:
        raise RuntimeError(label + '返回的方案格式损坏，请重试') from exc


def _mimo_key():
    key = os.environ.get('MIMO_API_KEY', '').strip()
    if not key:
        raise RuntimeError('PPT 新引擎未配置：请在 Railway Variables 添加 MIMO_API_KEY。不会再降级为旧版本地 PPT。')
    return key


def _mimo_chat(messages, *, web=False, temperature=0.55, timeout=100):
    """One MiMo turn.  Web search is a model tool, not fabricated source text."""
    payload = {
        'model': os.environ.get('MIMO_MODEL', 'mimo-v2.5').strip() or 'mimo-v2.5',
        'messages': messages,
        'temperature': temperature,
        'max_completion_tokens': 12000,
        'stream': False,
    }
    if web:
        # MiMo's Chat Completions API exposes this built-in tool.  Search results
        # are returned as response annotations and are collected below.
        payload['tools'] = [{'type': 'web_search'}]
    req = urllib.request.Request(
        'https://api.xiaomimimo.com/v1/chat/completions',
        data=json.dumps(payload).encode('utf-8'),
        headers={'api-key': _mimo_key(), 'Content-Type': 'application/json'},
    )
    try:
        result = json.loads(urllib.request.urlopen(req, timeout=timeout).read().decode('utf-8'))
    except urllib.error.HTTPError as exc:
        detail = exc.read().decode('utf-8', errors='ignore')[:500]
        raise RuntimeError('MiMo 请求失败（HTTP %s）：%s' % (exc.code, detail or exc.reason)) from exc
    choice = (result.get('choices') or [{}])[0]
    message = choice.get('message') or {}
    content = message.get('content') or ''
    if isinstance(content, list):
        content = ''.join(x.get('text', '') if isinstance(x, dict) else str(x) for x in content)
    if not str(content).strip():
        raise RuntimeError('MiMo 未返回内容')
    sources, seen = [], set()
    for annotation in message.get('annotations') or []:
        if not isinstance(annotation, dict):
            continue
        url = str(annotation.get('url') or annotation.get('url_citation', {}).get('url') or '').strip()
        title = str(annotation.get('title') or annotation.get('url_citation', {}).get('title') or '').strip()
        if url.startswith(('https://', 'http://')) and url not in seen:
            seen.add(url)
            sources.append({'title': title or urllib.parse.urlparse(url).netloc, 'url': url})
    return str(content).strip(), sources


def _detail_pages(detail):
    return {'短': '5–7', '中': '8–12', '长': '12–18'}.get(str(detail), '8–12')


def _deck_diagnostics(slides):
    warnings, titles = [], set()
    for index, slide in enumerate(slides, 1):
        title = str(slide.get('title') or '').strip()
        blocks = slide.get('blocks') or []
        chars = sum(len(str(x)) for x in blocks)
        if not title:
            warnings.append('第%d页没有标题' % index)
        if title in titles:
            warnings.append('第%d页标题重复：%s' % (index, title))
        titles.add(title)
        if len(title) > 32:
            warnings.append('第%d页标题过长（%d字）' % (index, len(title)))
        if chars > 620:
            warnings.append('第%d页正文过长（%d字），可能超框' % (index, chars))
        if not blocks:
            warnings.append('第%d页没有正文' % index)
    return warnings


def _normalise_agent_slides(plan):
    allowed = {'opening', 'narrative', 'explain', 'cards', 'timeline', 'scenario', 'activity', 'action', 'closing'}
    clean = []
    for item in plan.get('slides', []) if isinstance(plan, dict) else []:
        if not isinstance(item, dict):
            continue
        title = str(item.get('title') or '').strip()[:40]
        blocks = item.get('blocks') or item.get('content') or []
        if isinstance(blocks, str):
            blocks = [blocks]
        blocks = [str(x).strip() for x in blocks if str(x).strip()]
        if not title or not blocks:
            continue
        kind = str(item.get('layout') or item.get('type') or 'explain').lower()
        clean.append({
            'title': title,
            'layout': kind if kind in allowed else 'explain',
            'blocks': blocks[:7],
            'visual_query': str(item.get('visual_query') or '').strip()[:140],
            'speaker_note': str(item.get('speaker_note') or '').strip()[:500],
        })
    if len(clean) < 3:
        raise RuntimeError('MiMo 没有生成足够的页面，已停止导出')
    return clean


def _research_prompt(topic, grade, subject, extra, detail):
    audience = '中国中小学生'
    if grade:
        audience = grade + '学生'
    subject_hint = ('；学科/场景：' + subject) if subject else ''
    return '''你是资深中国教师与演示设计师。请先联网检索，再为“%s”设计一份真正可上课的 PPT。

对象：%s%s。教师补充：%s。目标篇幅：%s 页内容页，可由主题复杂度自然增减，不要为了凑页机械拆分。

先在心中完成研究与教学判断，再只输出一个 JSON 对象：
{"deck_title":"...","subtitle":"...","slides":[{"title":"...","layout":"opening|narrative|explain|cards|timeline|scenario|activity|action|closing","blocks":["..."],"visual_query":"适合真实图片搜索的具体中文词；不需要则空字符串","speaker_note":"教师讲述/追问提示"}]}

硬要求：
1. 必须联网，优先读取权威原始材料、专业机构资料、论文或正规新闻；检索只用于事实、数据、案例、图像线索，不能把搜索摘要拼成课件。
2. 用你的学科与教学知识组织内容。每页都要承担不同的教学作用：引入、解释、辨析、案例、活动、方法、迁移、收束可按需要选用，不能套固定八页模板。
3. 可以写完整、有节奏的段落，不要把所有内容压成口号；也不能把一句话拆成一页。每页文字须适合 16:9 投影展示。
4. 未被可靠来源支持时，不写精确数字、具体新闻细节、法规条款或研究结论；不要虚构教师、学校、学生或采访。
5. visual_query 只能用于真正能帮助理解的页面，写“真实照片/现场图/示意图”的检索词；禁止要求图片内含文字、海报、编号或水印。
6. 不要在正文中塞 URL、参考文献或“据搜索结果”。来源由系统在最后统一列出。''' % (topic, audience, subject_hint, extra or '无', _detail_pages(detail))


def build_agent_deck(body):
    """Research → content plan → machine-readable QA → revision. No old fallback."""
    topic = str(body.get('topic') or body.get('title') or '').strip()
    if not topic:
        raise RuntimeError('PPT主题不能为空')
    detail = str(body.get('detail') or body.get('wordCount') or '中')
    grade = str(body.get('grade') or '').strip()
    subject = str(body.get('subject') or '').strip()
    extra = str(body.get('userDetail') or body.get('harvest') or '').strip()[:1600]
    trace = bool(body.get('debug'))
    print('[PPT AGENT] 01 research-and-plan topic=%s detail=%s' % (topic, detail))
    raw, sources = _mimo_chat([
        {'role': 'system', 'content': '你产出的是给教师直接使用的演示方案。宁可少写未经证实的事实，也不要编造。'},
        {'role': 'user', 'content': _research_prompt(topic, grade, subject, extra, detail)},
    ], web=True, temperature=0.7, timeout=115)
    if len(sources) < 2:
        raise RuntimeError('联网搜索未返回至少 2 个可追溯来源；为避免假资料，已停止生成')
    initial = _normalise_agent_slides(_extract_json_object(raw, 'MiMo'))
    diagnostics = _deck_diagnostics(initial)
    print('[PPT AGENT] 02 research sources=%d, initial slides=%d, diagnostics=%s' % (
        len(sources), len(initial), '；'.join(diagnostics) if diagnostics else '通过'))
    revision_prompt = '''下面是一份已经完成联网研究的教师 PPT 草案。请作为苛刻的教研员和演示设计师修订它。
主题：%s
对象：%s
机器检测：%s

只输出与原结构完全相同的 JSON 对象。保留有价值的完整表述；删掉空话、重复、没有依据的精确事实；必要时合并或拆分页面。确保每页在投影上能读完，并让版式类型有节奏变化。不要新编来源、案例或数据。\n\n草案：%s''' % (
        topic, grade or '中小学生', '；'.join(diagnostics) or '无明显结构问题',
        json.dumps({'deck_title': topic, 'slides': initial}, ensure_ascii=False))
    revised_raw, _ = _mimo_chat([
        {'role': 'system', 'content': '只返回 JSON，不写解释。'},
        {'role': 'user', 'content': revision_prompt},
    ], web=False, temperature=0.35, timeout=100)
    revised = _extract_json_object(revised_raw, 'MiMo修订')
    slides = _normalise_agent_slides(revised)
    final_diagnostics = _deck_diagnostics(slides)
    print('[PPT AGENT] 03 revision slides=%d, diagnostics=%s' % (
        len(slides), '；'.join(final_diagnostics) if final_diagnostics else '通过'))
    deck = {
        'title': str(revised.get('deck_title') or topic).strip()[:60],
        'subtitle': str(revised.get('subtitle') or '').strip()[:100],
        'theme': detect(topic),
        'slides': slides,
        'sources': sources[:8],
        'debug': ({'research_sources': sources, 'initial_plan': initial, 'initial_diagnostics': diagnostics,
                   'revised_plan': slides, 'final_diagnostics': final_diagnostics} if trace else None),
    }
    return deck


def build_agent_research(topic, detail='中'):
    """PPT research endpoint uses the same research engine and source guarantee."""
    raw, sources = _mimo_chat([
        {'role': 'system', 'content': '你是严谨的中国教师备课研究员。'},
        {'role': 'user', 'content': '请联网检索“%s”的课堂资料。只写可验证事实、可用案例方向和教学建议；标出不宜当作事实使用的内容。篇幅：%s。' % (topic, detail)},
    ], web=True, temperature=0.25, timeout=90)
    if len(sources) < 2:
        raise RuntimeError('联网搜索没有返回足够可核验来源')
    return {'research_note': raw[:5000], 'sources': sources[:8], 'assets': []}


def _safe_download_image(url):
    request = urllib.request.Request(url, headers={'User-Agent': 'Mozilla/5.0 (compatible; ShizhuAI/1.0)'})
    with urllib.request.urlopen(request, timeout=18) as response:
        blob = response.read(8 * 1024 * 1024 + 1)
        content_type = response.headers.get_content_type()
    if len(blob) > 8 * 1024 * 1024:
        raise RuntimeError('图片过大')
    if not content_type.startswith('image/'):
        raise RuntimeError('链接不是图片')
    suffix = '.png' if 'png' in content_type else '.jpg'
    path = os.path.join(tempfile.mkdtemp(prefix='ppt_visual_'), 'asset' + suffix)
    with open(path, 'wb') as fh:
        fh.write(blob)
    return path


def _attach_research_visuals(deck, max_images=4):
    """Every planned visual gets one attempt. Failed external downloads never invent AI art."""
    image_key = os.environ.get('AI_API_KEY', '').strip()
    if not image_key:
        print('[PPT AGENT] 04 visual search skipped: AI_API_KEY unavailable')
        return
    completed = 0
    for index, slide in enumerate(deck['slides'], 1):
        if completed >= max_images:
            break
        query = slide.get('visual_query', '')
        if not query:
            continue
        try:
            url = search_image_url(query, image_key, timeout=20)
            slide['_image_path'] = _safe_download_image(url)
            completed += 1
            print('[PPT AGENT] 04 visual %d attached: %s' % (index, query))
        except Exception as exc:
            print('[PPT AGENT] 04 visual %d skipped: %s' % (index, str(exc)[:160]))
    print('[PPT AGENT] 04 visuals complete=%d' % completed)


def _add_textbox(slide, x, y, w, h, text, size, color, *, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Pt(4); frame.margin_right = Pt(4)
    para = frame.paragraphs[0]
    para.text = str(text)
    para.font.size = Pt(size); para.font.color.rgb = color; para.font.bold = bold; para.alignment = align
    return box


def _add_card(slide, x, y, w, h, text, fill, color, size=17):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid(); card.fill.fore_color.rgb = fill; card.line.fill.background()
    frame = card.text_frame; frame.word_wrap = True
    frame.margin_left = Pt(12); frame.margin_right = Pt(12); frame.margin_top = Pt(8)
    p = frame.paragraphs[0]; p.text = str(text); p.font.size = Pt(size); p.font.color.rgb = color
    return card


def render_agent_deck(deck):
    """A small deliberate design system. Content controls layout; renderer controls overflow."""
    _attach_research_visuals(deck)
    theme = THEME.get(deck.get('theme'), THEME['default'])
    dark, primary, accent, light = [rgb(x) for x in theme]
    white, text = RGBColor(255, 255, 255), RGBColor(45, 50, 60)
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)

    all_slides = [{'title': deck['title'], 'layout': 'cover', 'blocks': [deck.get('subtitle', '')]}] + deck['slides']
    all_slides.append({'title': '资料来源与延伸阅读', 'layout': 'sources', 'blocks': deck.get('sources', [])})
    all_slides.append({'title': '谢谢聆听', 'layout': 'end', 'blocks': []})

    for index, item in enumerate(all_slides):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        kind = item.get('layout', 'explain')
        blocks = item.get('blocks') or []
        if kind == 'cover':
            slide.background.fill.solid(); slide.background.fill.fore_color.rgb = dark
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(.9), Inches(1.1), Inches(.14), Inches(4.9))
            bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
            _add_textbox(slide, 1.35, 2.1, 10.5, 1.6, item['title'], 42, white, bold=True)
            if blocks and blocks[0]:
                _add_textbox(slide, 1.4, 4.0, 9.8, .7, blocks[0], 20, light)
            continue
        if kind == 'end':
            slide.background.fill.solid(); slide.background.fill.fore_color.rgb = light
            _add_textbox(slide, 1.2, 2.7, 10.9, 1, item['title'], 40, dark, bold=True, align=PP_ALIGN.CENTER)
            continue
        # Header and page marker for every work page.
        _add_textbox(slide, .8, .45, 11.5, .65, item['title'], 29, dark, bold=True)
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(.82), Inches(1.2), Inches(1.15), Pt(4))
        line.fill.solid(); line.fill.fore_color.rgb = accent; line.line.fill.background()
        _add_textbox(slide, 11.85, 6.92, .65, .25, '%d/%d' % (index + 1, len(all_slides)), 9, primary, align=PP_ALIGN.RIGHT)

        if kind == 'sources':
            y = 1.65
            for source in blocks[:8]:
                domain = urllib.parse.urlparse(source.get('url', '')).netloc
                _add_card(slide, 1.0, y, 11.1, .52, '%s  ·  %s' % (source.get('title', domain), domain), light, text, 13)
                y += .62
            continue

        image_path = item.get('_image_path')
        if kind in ('narrative', 'scenario', 'opening'):
            x, w = (5.15, 6.75) if image_path else (1.05, 11.1)
            panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.6), Inches(w), Inches(4.8))
            panel.fill.solid(); panel.fill.fore_color.rgb = light; panel.line.fill.background()
            joined = '\n\n'.join(blocks[:2])
            font = 20 if len(joined) < 280 else 17 if len(joined) < 430 else 15
            _add_textbox(slide, x + .25, 1.9, w - .5, 4.15, joined, font, text)
            if image_path:
                slide.shapes.add_picture(image_path, Inches(1.0), Inches(1.75), Inches(3.7), Inches(3.7))
            continue

        if kind in ('cards', 'action', 'activity', 'timeline'):
            count = min(max(len(blocks), 1), 6)
            cols = 3 if count > 4 else 2 if count > 1 else 1
            rows = (count + cols - 1) // cols
            card_w = 11.25 / cols - .16
            card_h = min(2.15, 4.75 / rows)
            for pos, block in enumerate(blocks[:6]):
                col, row = pos % cols, pos // cols
                x, y = 1.0 + col * (11.25 / cols), 1.6 + row * (4.85 / rows)
                _add_card(slide, x, y, card_w, card_h, block, light, text, 16 if len(block) < 85 else 14)
                dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + .16), Inches(y + .17), Inches(.28), Inches(.28))
                dot.fill.solid(); dot.fill.fore_color.rgb = accent; dot.line.fill.background()
            continue

        # Explain is intentionally dense enough for real teaching, but resizes
        # before overflow rather than forcing arbitrary page splits.
        x, w = (1.0, 6.7) if image_path else (1.0, 11.25)
        y = 1.55
        for block in blocks[:6]:
            height = .72 if len(block) < 75 else 1.02 if len(block) < 150 else 1.32
            _add_card(slide, x, y, w, height, block, light, text, 16 if len(block) < 140 else 14)
            y += height + .18
            if y > 6.3:
                break
        if image_path:
            slide.shapes.add_picture(image_path, Inches(8.1), Inches(1.8), Inches(3.55), Inches(3.55))

    output = io.BytesIO(); prs.save(output)
    print('[PPT AGENT] 05 rendered slides=%d bytes=%d' % (len(all_slides), output.tell()))
    return output.getvalue()


def _hex_color(value, fallback):
    value = str(value or '').strip()
    if len(value) == 7 and value.startswith('#'):
        try:
            return rgb(value)
        except ValueError:
            pass
    return fallback


def _freeform_prompt(topic, grade, subject, extra):
    audience = (grade + '学生') if grade else '中小学生'
    subject_line = ('；相关学科：' + subject) if subject else ''
    return '''你是一位有经验的中国中小学教师，也是一位优秀的演示设计师。请为教师直接生成一份可使用的 PPT，可以是课堂教学课件、主题班会、家长会、工作汇报或其他教育教学场景。根据主题自主联网查阅资料，并结合你的知识完成内容与视觉设计。

主题：%s
对象：%s%s
教师补充：%s

请使用你自己的判断决定页数、结构、文字篇幅、讲述顺序和每一页的视觉设计。让内容真正适合课堂或班会使用、好讲也好懂。需要图片时，写出合适的真实图片搜索词。

只输出一个 JSON 对象，不要 Markdown 或解释。它是 13.333×7.5 英寸画布的自由演示设计：
{
 "title":"总标题",
 "slides":[
   {"background":"#FFFFFF", "elements":[
      {"kind":"shape","shape":"rect|round_rect|line|ellipse","x":0,"y":0,"w":13.333,"h":7.5,"fill":"#FFFFFF","line":"#FFFFFF"},
      {"kind":"text","x":1,"y":1,"w":8,"h":1,"text":"文字","font_size":28,"color":"#1F2937","bold":true,"align":"left|center|right"},
      {"kind":"image","x":8,"y":1.5,"w":4,"h":3,"query":"适合检索真实图片的中文词"}
   ]}
 ]
}
坐标和尺寸直接使用英寸。先放背景/形状，再放文字和图片；图片元素按页面需要添加。''' % (topic, audience, subject_line, extra or '无')


def _normalise_freeform_deck(plan):
    if not isinstance(plan, dict) or not isinstance(plan.get('slides'), list):
        raise RuntimeError('MiMo 没有返回 PPT 页面')
    slides = []
    for raw_slide in plan['slides'][:30]:
        if not isinstance(raw_slide, dict):
            continue
        elements = []
        for raw in raw_slide.get('elements', [])[:45]:
            if not isinstance(raw, dict):
                continue
            kind = str(raw.get('kind') or '').lower()
            if kind not in {'text', 'shape', 'image'}:
                continue
            try:
                x, y = float(raw.get('x', 0)), float(raw.get('y', 0))
                w, h = float(raw.get('w', 1)), float(raw.get('h', 1))
            except (TypeError, ValueError):
                continue
            # This only keeps elements on the physical slide; it does not choose
            # their layout or rewrite their content.
            x, y = max(0, min(x, 13.2)), max(0, min(y, 7.4))
            w, h = max(.08, min(w, 13.333 - x)), max(.08, min(h, 7.5 - y))
            item = {'kind': kind, 'x': x, 'y': y, 'w': w, 'h': h}
            if kind == 'text':
                text_value = str(raw.get('text') or '').strip()
                if not text_value:
                    continue
                item.update({'text': text_value, 'font_size': max(8, min(float(raw.get('font_size', 18)), 52)),
                             'color': str(raw.get('color') or '#1F2937'), 'bold': bool(raw.get('bold')),
                             'align': str(raw.get('align') or 'left').lower()})
            elif kind == 'shape':
                item.update({'shape': str(raw.get('shape') or 'rect').lower(),
                             'fill': str(raw.get('fill') or '#FFFFFF'), 'line': str(raw.get('line') or raw.get('fill') or '#FFFFFF')})
            else:
                query = str(raw.get('query') or '').strip()
                if not query:
                    continue
                item['query'] = query[:180]
            elements.append(item)
        if elements:
            slides.append({'background': str(raw_slide.get('background') or '#FFFFFF'), 'elements': elements})
    if len(slides) < 2:
        raise RuntimeError('MiMo 未生成足够页面')
    return {'title': str(plan.get('title') or '').strip()[:80], 'slides': slides}


def build_freeform_deck(body):
    topic = str(body.get('topic') or body.get('title') or '').strip()
    if not topic:
        raise RuntimeError('PPT主题不能为空')
    print('[PPT FREEFORM] 01 research-and-design topic=%s' % topic)
    raw, sources = _mimo_chat([
        {'role': 'system', 'content': '你要自由设计一份可直接使用的中小学教师 PPT。'},
        {'role': 'user', 'content': _freeform_prompt(topic, str(body.get('grade') or '').strip(),
                                                       str(body.get('subject') or '').strip(),
                                                       str(body.get('userDetail') or body.get('harvest') or '').strip()[:1600])},
    ], web=True, temperature=0.72, timeout=140)
    deck = _normalise_freeform_deck(_extract_json_object(raw, 'MiMo'))
    if not deck['title']:
        deck['title'] = topic
    deck['sources'] = sources[:10]
    deck['debug'] = {'research_sources': sources, 'design': deck} if body.get('debug') else None
    print('[PPT FREEFORM] 02 sources=%d slides=%d' % (len(sources), len(deck['slides'])))
    return deck


def render_freeform_deck(deck):
    """Render the model's free canvas. No fixed page count, outline, or layout."""
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    default_text = RGBColor(31, 41, 55)
    for page_number, spec in enumerate(deck['slides'], 1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid(); slide.background.fill.fore_color.rgb = _hex_color(spec.get('background'), RGBColor(255, 255, 255))
        for element in spec['elements']:
            x, y, w, h = (Inches(element['x']), Inches(element['y']), Inches(element['w']), Inches(element['h']))
            if element['kind'] == 'shape':
                shape_map = {'round_rect': MSO_SHAPE.ROUNDED_RECTANGLE, 'ellipse': MSO_SHAPE.OVAL,
                             'line': MSO_SHAPE.RECTANGLE, 'rect': MSO_SHAPE.RECTANGLE}
                shape = slide.shapes.add_shape(shape_map.get(element.get('shape'), MSO_SHAPE.RECTANGLE), x, y, w, h)
                shape.fill.solid(); shape.fill.fore_color.rgb = _hex_color(element.get('fill'), RGBColor(255, 255, 255))
                shape.line.color.rgb = _hex_color(element.get('line'), RGBColor(255, 255, 255))
            elif element['kind'] == 'text':
                align_map = {'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT, 'left': PP_ALIGN.LEFT}
                _add_textbox(slide, element['x'], element['y'], element['w'], element['h'], element['text'],
                             element['font_size'], _hex_color(element.get('color'), default_text),
                             bold=element.get('bold', False), align=align_map.get(element.get('align'), PP_ALIGN.LEFT))
            elif element['kind'] == 'image':
                # The deck remains usable if a legitimate public image cannot be retrieved.
                key = os.environ.get('AI_API_KEY', '').strip()
                if not key:
                    continue
                try:
                    image_url = search_image_url(element['query'], key, timeout=20)
                    slide.shapes.add_picture(_safe_download_image(image_url), x, y, w, h)
                    print('[PPT FREEFORM] image attached page=%d query=%s' % (page_number, element['query']))
                except Exception as exc:
                    print('[PPT FREEFORM] image skipped page=%d: %s' % (page_number, str(exc)[:160]))
    output = io.BytesIO(); prs.save(output)
    print('[PPT FREEFORM] 03 exported slides=%d bytes=%d' % (len(deck['slides']), output.tell()))
    return output.getvalue()


def _native_ppt_prompt(topic, grade, subject, extra, detail):
    """A constrained story-board is more reliable than asking a model for coordinates."""
    audience = (grade + '学生') if grade else '中小学生'
    subject_line = ('；相关学科：' + subject) if subject else ''
    pages = {'短': '5-6', '中': '7-9', '长': '10-12'}.get(str(detail), '7-9')
    return '''你是一位中国中小学教师和专业演示设计师。请先联网核实需要核实的事实，再为主题设计一份能直接用于课堂、班会或家长会的 PPT 故事板。

主题：%s
对象：%s%s
教师补充：%s
内容页目标：%s 页。封面、资料来源、结束页由程序生成，不要输出。

只输出一个合法 JSON 对象，不要 Markdown 或解释：
{"title":"总标题","subtitle":"简短副标题","slides":[{"title":"页面标题","layout":"story|explain|cards|steps|timeline|compare|scenario|action","blocks":["投影上可读的完整文字"],"visual_query":"需要真实配图时填写具体中文检索词，否则为空字符串","speaker_note":"给教师的讲述或追问提示"}]}

硬性质量标准：
1. 每页只解决一个教学问题，结构随主题自然推进，不套固定八页模板；整套至少使用 4 种不同 layout。
2. story/scenario 可用 1-2 段完整、通顺的文字；explain 用 2-5 条有解释力的要点；cards/action 用 2-6 条可执行内容；steps/timeline 用 3-5 步；compare 的 blocks 前半为左侧、后半为右侧。不要把同一句话拆成续页。
3. 事实、数据、政策、真实事件必须经联网资料支撑；不确定时使用课堂情境或条件表达，绝不虚构学校、学生、采访或来源。
4. 不要在 blocks 中放 URL、参考文献、页码、项目符号、编号或“资料显示”；来源由系统统一生成。
5. visual_query 仅在一张真实图片明显帮助理解时填写，描述具体、无文字、无水印的真实照片或示意图。每页至多一张图，整套不必强行配图。
6. 语言符合 %s 的理解水平，少用口号，多给解释、判断方法、可讨论的情境或可落实的行动。''' % (topic, audience, subject_line, extra or '无', pages, audience)


def _normalise_native_deck(plan, topic, detail):
    """Validate model output before it reaches the editable native renderer."""
    if not isinstance(plan, dict) or not isinstance(plan.get('slides'), list):
        raise RuntimeError('PPT方案没有返回页面结构')
    allowed = {'story', 'explain', 'cards', 'steps', 'timeline', 'compare', 'scenario', 'action'}
    max_pages = {'短': 6, '中': 9, '长': 12}.get(str(detail), 9)
    slides = []
    for raw in plan['slides'][:max_pages]:
        if not isinstance(raw, dict):
            continue
        title = str(raw.get('title') or '').strip()[:40]
        blocks = raw.get('blocks') or raw.get('content') or []
        if isinstance(blocks, str):
            blocks = [blocks]
        blocks = [str(item).strip()[:260] for item in blocks if str(item).strip()]
        if not title or not blocks:
            continue
        layout = str(raw.get('layout') or raw.get('type') or 'explain').lower()
        slides.append({
            'title': title,
            'layout': layout if layout in allowed else 'explain',
            'blocks': blocks[:6],
            'visual_query': str(raw.get('visual_query') or '').strip()[:140],
            'speaker_note': str(raw.get('speaker_note') or '').strip()[:700],
        })
    if len(slides) < 3:
        raise RuntimeError('PPT内容页不足，未导出')
    # A repeated layout is permissible for short decks, but not a whole deck.
    layouts = {slide['layout'] for slide in slides}
    if len(slides) >= 5 and len(layouts) < 3:
        raise RuntimeError('PPT版式变化不足，已停止导出以避免生成模板化页面')
    return {
        'title': str(plan.get('title') or topic).strip()[:60] or topic,
        'subtitle': str(plan.get('subtitle') or '').strip()[:100],
        'slides': slides,
    }


def build_native_deck(body):
    """Use the existing Alibaba Qwen key to research and write the native storyboard."""
    topic = str(body.get('topic') or body.get('title') or '').strip()
    if not topic:
        raise RuntimeError('PPT主题不能为空')
    detail = str(body.get('detail') or body.get('wordCount') or '中')
    grade = str(body.get('grade') or '').strip()
    subject = str(body.get('subject') or '').strip()
    extra = str(body.get('userDetail') or body.get('harvest') or '').strip()[:1600]
    api_key = os.environ.get('AI_API_KEY', '').strip()
    if not api_key:
        raise RuntimeError('原生PPT服务未配置 AI_API_KEY，请在 Railway 设置阿里百炼 API Key')
    print('[PPT NATIVE] 01 qwen-research-and-storyboard topic=%s detail=%s' % (topic, detail))
    # _qwen_search_storyboard uses Alibaba Qwen Plus with forced web search and
    # returns only citations supplied by the actual search response. No MiMo key
    # or Alibaba official-PPT mixed-billing service is involved.
    raw, sources = _qwen_search_storyboard(
        api_key,
        _native_ppt_prompt(topic, grade, subject, extra, detail),
        timeout=135,
    )
    deck = _normalise_native_deck(_extract_json_object(raw, 'MiMo'), topic, detail)
    deck['theme'] = str(body.get('theme') or detect(topic))
    deck['sources'] = sources[:6]
    if body.get('debug'):
        deck['debug'] = {'research_sources': sources, 'storyboard': deck['slides']}
    print('[PPT NATIVE] 02 sources=%d slides=%d layouts=%s' % (
        len(sources), len(deck['slides']), ','.join(sorted({s['layout'] for s in deck['slides']}))))
    return deck


def render_native_deck(deck):
    """Use PptxGenJS so all text, shapes and images remain editable in PowerPoint."""
    _attach_research_visuals(deck, max_images=4)
    renderer = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'ppt_native_renderer.js')
    if not os.path.isfile(renderer):
        raise RuntimeError('原生PPT渲染器文件缺失')
    workdir = tempfile.mkdtemp(prefix='native_ppt_')
    input_path = os.path.join(workdir, 'deck.json')
    output_path = os.path.join(workdir, 'deck.pptx')
    with open(input_path, 'w', encoding='utf-8') as fh:
        json.dump(deck, fh, ensure_ascii=False)
    try:
        completed = subprocess.run(
            ['node', renderer, input_path, output_path],
            capture_output=True, text=True, timeout=90, check=False,
        )
    except FileNotFoundError as err:
        raise RuntimeError('服务器未安装 Node.js，无法生成原生可编辑 PPT') from err
    except subprocess.TimeoutExpired as err:
        raise RuntimeError('原生PPT排版超时') from err
    if completed.returncode != 0:
        raise RuntimeError('原生PPT渲染失败：' + (completed.stderr or completed.stdout or '未知错误')[:500])
    if not os.path.isfile(output_path) or os.path.getsize(output_path) < 5000:
        raise RuntimeError('原生PPT导出文件异常')
    with open(output_path, 'rb') as fh:
        result = fh.read()
    print('[PPT NATIVE] 03 rendered bytes=%d' % len(result))
    return result


def gen_docx(data):
    """生成Word文档(python-docx)，带【】小标题识别、字体、字号、缩进"""
    from docx import Document
    from docx.shared import Pt, Cm, Emu
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    
    doc = Document()
    
    # 默认样式：宋体 小四 1.3倍行距
    style = doc.styles['Normal']
    style.font.name = 'SimSun'
    style.font.size = Pt(12)
    style.element.rPr.rFonts.set(qn('w:eastAsia'), '宋体')
    pf = style.paragraph_format
    pf.line_spacing = 1.3
    
    title = data.get('title', '文档')
    doc_title = data.get('docTitle', '') or title
    doc_subtitle = data.get('docSubtitle', '')
    
    # 大标题：小二 居中 加粗（纯文本，无蓝色主题）
    tp = doc.add_paragraph()
    tp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    tr = tp.add_run(doc_title)
    tr.bold = True
    tr.font.size = Pt(18)
    tr.font.name = 'SimSun'
    tr.font.color.rgb = None
    tr.element.rPr.rFonts.set(qn('w:eastAsia'), '宋体')
    
    if doc_subtitle:
        sp = doc.add_paragraph()
        sp.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sr = sp.add_run(doc_subtitle)
        sr.bold = True
        sr.font.size = Pt(14)
        sr.font.name = 'SimSun'
        sr.element.rPr.rFonts.set(qn('w:eastAsia'), '宋体')
    
    content = data.get('content', '')
    subtitle_counter = 0
    
    if content:
        for para in content.split('\n'):
            para = para.strip()
            if not para:
                continue
            try:
                # 【小标题】：加粗大号、自动编号、不缩进
                import re
                bracket_match = re.search(r'【(.+?)】', para)
                if bracket_match:
                    sub_text = bracket_match.group(1)
                    subtitle_counter += 1
                    if not re.match(r'^[一二三四五六七八九十\d]', sub_text):
                        sub_text = str(subtitle_counter) + '. ' + sub_text
                    # 去掉【】部分
                    para = re.sub(r'【.+?】', '', para).strip()
                    sp = doc.add_paragraph()
                    sp.paragraph_format.space_before = Pt(12)
                    sp.paragraph_format.space_after = Pt(6)
                    sp.paragraph_format.line_spacing = 1.3
                    sr = sp.add_run(sub_text)
                    sr.bold = True
                    sr.font.size = Pt(15)
                    sr.font.name = 'SimSun'
                    sr.element.rPr.rFonts.set(qn('w:eastAsia'), '宋体')
                    if para:
                        dp = doc.add_paragraph(para)
                        dp.paragraph_format.first_line_indent = Pt(24)
                        dp.paragraph_format.line_spacing = 1.3
                        dp.paragraph_format.space_after = Pt(6)
                    continue
                
                if para.startswith('# ') or para.startswith('## '):
                    hp = doc.add_paragraph()
                    hp.paragraph_format.line_spacing = 1.3
                    hr = hp.add_run(para[para.index(' ')+1:])
                    hr.bold = True
                    hr.font.size = Pt(15)
                    hr.font.name = 'SimSun'
                    hr.font.color.rgb = None
                    hr.element.rPr.rFonts.set(qn('w:eastAsia'), '宋体')
                elif para.startswith('**') and para.endswith('**'):
                    d = doc.add_paragraph()
                    d.paragraph_format.line_spacing = 1.3
                    r = d.add_run(para[2:-2])
                    r.bold = True
                    r.font.size = Pt(12)
                    r.font.name = 'SimSun'
                    r.element.rPr.rFonts.set(qn('w:eastAsia'), '宋体')
                elif para.startswith('- ') or para.startswith('* '):
                    li = doc.add_paragraph(para[2:], style='List Bullet')
                    li.paragraph_format.line_spacing = 1.3
                else:
                    dp = doc.add_paragraph(para)
                    dp.paragraph_format.first_line_indent = Pt(24)
                    dp.paragraph_format.line_spacing = 1.3
                    dp.paragraph_format.space_after = Pt(6)
            except Exception:
                dp = doc.add_paragraph(para)
                dp.paragraph_format.line_spacing = 1.3
                dp.paragraph_format.space_after = Pt(6)
    
    buf = __import__('io').BytesIO()
    doc.save(buf)
    return buf.getvalue()

def parse_file(data):
    """解析文档/PDF/Excel，返回文本内容"""
    import tempfile, os
    content = data.get('content', '')  # base64 encoded file
    filename = data.get('filename', '')
    if not content:
        return '（无文件内容）'
    
    import base64
    raw = base64.b64decode(content)
    ext = filename.lower().split('.')[-1] if '.' in filename else ''
    
    try:
        if ext in ('txt', 'csv', 'json', 'md', 'xml'):
            return raw.decode('utf-8')
        elif ext in ('jpg', 'jpeg', 'png', 'gif', 'bmp', 'webp'):
            api_key = os.environ.get('AI_API_KEY', '').strip()
            if not api_key:
                raise RuntimeError('图片文字识别未配置：请在 Railway Variables 中设置 AI_API_KEY')
            import urllib.request, urllib.error, base64
            img_b64 = base64.b64encode(raw).decode('ascii')
            # Auto-detect image format
            if raw[:4] == b'\x89PNG':
                mime_type = 'png'
            elif raw[:3] == b'\xff\xd8\xff':
                mime_type = 'jpeg'
            elif raw[:4] in (b'GIF8',):
                mime_type = 'gif'
            elif raw[:2] == b'BM':
                mime_type = 'bmp'
            else:
                mime_type = 'jpeg'
            ocr_data = json.dumps({
                'model': 'qwen3-vl-flash',
                'input': {
                    'messages': [{'role': 'user', 'content': [
                        {'image': 'data:image/' + mime_type + ';base64,' + img_b64},
                        {'text': '请提取这张图片中的所有文字内容，直接输出文字，不要额外说明'}
                    ]}]
                }
            }).encode()
            req = urllib.request.Request(
                'https://ws-5ol6m5p8f4hikz1a.cn-beijing.maas.aliyuncs.com/api/v1/services/aigc/multimodal-generation/generation',
                data=ocr_data,
                headers={'Authorization': 'Bearer ' + api_key, 'Content-Type': 'application/json'}
            )
            # 截图 OCR 正常应在数十秒内完成；超时要返回明确错误，不能让小程序一直转圈。
            response = urllib.request.urlopen(req, timeout=45)
            result = json.loads(response.read().decode('utf-8'))
            text = result['output']['choices'][0]['message']['content']
            if isinstance(text, list):
                text = ''.join(
                    part.get('text', '') if isinstance(part, dict) else str(part)
                    for part in text
                )
            return '【图片：' + filename + '】\n' + str(text).strip()
        elif ext == 'pdf':
            import pdfplumber
            import io
            text = ''
            with pdfplumber.open(io.BytesIO(raw)) as pdf:
                for page in pdf.pages:
                    text += page.extract_text() or ''
            return text[:50000]
        elif ext in ('docx', 'doc'):
            import docx
            import io
            doc = docx.Document(io.BytesIO(raw))
            return '\n'.join([p.text for p in doc.paragraphs])[:50000]
        elif ext in ('pptx', 'ppt'):
            from pptx import Presentation
            import io
            prs = Presentation(io.BytesIO(raw))
            text = ''
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text:
                        text += shape.text + '\n'
            return text[:50000] if text else '（PPT文件无可提取文字）'
        elif ext in ('xlsx', 'xls'):
            import openpyxl
            import io
            wb = openpyxl.load_workbook(io.BytesIO(raw))
            rows = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                rows.append(f'【工作表：{sheet}】')
                for row in ws.iter_rows(values_only=True):
                    rows.append('\t'.join([str(c) if c is not None else '' for c in row]))
            return '\n'.join(rows)[:50000]
        else:
            return f'不支持的文件格式：.{ext}'
    except Exception as e:
        # 让 HTTP 层返回 code=-1。不能把异常伪装成文件文字，
        # 否则前端会继续把它交给 DeepSeek 做“文件总结”。
        raise RuntimeError(f'文件解析失败：{str(e)}') from e

def edit_excel(data):
    """根据AI指令修改Excel文件"""
    import openpyxl, re
    file_b64 = data.get('file', '')
    instructions = data.get('instructions', '')
    if not file_b64:
        return None, '未提供Excel文件'
    
    try:
        raw = base64.b64decode(file_b64)
        wb = openpyxl.load_workbook(io.BytesIO(raw))
    except:
        wb = openpyxl.Workbook()
    
    ws = wb.active
    
    # 解析简单指令：修改单元格
    # 格式："设置 A1=张三" 或 "修改 B2=95"
    for line in instructions.split('\n'):
        line = line.strip()
        m = re.match(r'设置\s*([A-Z]+)(\d+)\s*=\s*(.+)', line)
        if m:
            cell = m.group(1) + m.group(2)
            ws[cell] = m.group(3).strip()
    
    buf = io.BytesIO()
    wb.save(buf)
    return base64.b64encode(buf.getvalue()).decode('ascii'), None

# === HTTP Server ===
def make_handler():
    class Handler(http.server.BaseHTTPRequestHandler):
        def _set_cors(self):
            self.send_header('Access-Control-Allow-Origin', '*')
            self.send_header('Access-Control-Allow-Methods', 'POST, OPTIONS')
            self.send_header('Access-Control-Allow-Headers', 'Content-Type')

        def do_OPTIONS(self):
            self.send_response(204)
            self._set_cors()
            self.end_headers()

        def do_GET(self):
            if self.path.startswith('/dl/'):
                file_id = self.path.split('/')[-1]
                for f in os.listdir(DL_DIR):
                    if f == file_id or f.startswith(file_id + '.'):
                        fpath = os.path.join(DL_DIR, f)
                        try:
                            with open(fpath, 'rb') as fh:
                                data = fh.read()
                            os.remove(fpath)
                            ext = f.rsplit('.', 1)[-1].lower() if '.' in f else 'bin'
                            content_types = {
                                'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                                'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                                'xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                            }
                            self.send_response(200)
                            self.send_header('Content-Type', content_types.get(ext, 'application/octet-stream'))
                            self.send_header('Content-Disposition', 'attachment; filename="szhuAI.' + ext + '"')
                            self._set_cors()
                            self.end_headers()
                            self.wfile.write(data)
                            return
                        except: pass
                self.send_response(404)
                self._set_cors()
                self.end_headers()
                return
            self.send_response(200)
            self.send_header('Content-Type', 'text/plain; charset=utf-8')
            self._set_cors()
            self.end_headers()
            self.wfile.write(('师助AI PPT API is running. pipeline=' + PPT_PIPELINE_VERSION).encode('utf-8'))

        def do_POST(self):
            try:
                length = int(self.headers.get('Content-Length', 0))
                body = json.loads(self.rfile.read(length))
                doc_type = body.get('type', 'ppt')
                
                if doc_type == 'create-xlsx':
                    import openpyxl
                    wb = openpyxl.Workbook()
                    ws = wb.active
                    ws.title = '学生名单'
                    names = body.get('names', [])
                    ws.cell(1, 1, '序号')
                    ws.cell(1, 2, '姓名')
                    for i, name in enumerate(names):
                        ws.cell(i+2, 1, i+1)
                        ws.cell(i+2, 2, name)
                    buf = io.BytesIO()
                    wb.save(buf)
                    # Save to temp file and return download URL
                    fid = str(uuid.uuid4())
                    fpath = os.path.join(DL_DIR, fid + '.xlsx')
                    with open(fpath, 'wb') as fh:
                        fh.write(buf.getvalue())
                    resp = json.dumps({'code': 0, 'url': '/dl/' + fid})
                    print(f'Created xlsx: {len(names)} names, url=/dl/{fid}')
                elif doc_type == 'seat-export':
                    import openpyxl
                    wb = openpyxl.Workbook()
                    ws = wb.active
                    ws.title = '座位表'
                    grid = body.get('grid', [])
                    col_types = body.get('cols', [])
                    ws.cell(1, 1, '排')
                    for ci, ct in enumerate(col_types):
                        if ct == 'seat':
                            ws.cell(1, ci+2, f'列{ci+1}')
                        else:
                            ws.cell(1, ci+2, '过道')
                    for ri, row in enumerate(grid):
                        ws.cell(ri+2, 1, f'第{ri+1}排')
                        for ci, cell in enumerate(row):
                            if isinstance(cell, dict) and cell.get('student'):
                                ws.cell(ri+2, ci+2, cell['student'])
                    buf = io.BytesIO()
                    wb.save(buf)
                    # Save to temp file and return download URL
                    fid = str(uuid.uuid4())
                    fpath = os.path.join(DL_DIR, fid + '.xlsx')
                    with open(fpath, 'wb') as fh:
                        fh.write(buf.getvalue())
                    resp = json.dumps({'code': 0, 'url': '/dl/' + fid})
                    print(f'Seat export: {len(grid)} rows, url=/dl/{fid}')
                elif doc_type == 'doc':
                    print(f'Generating DOC: {body.get("title","")}')
                    docx_bytes = gen_docx(body)
                    b64 = base64.b64encode(docx_bytes).decode('ascii')
                    resp = json.dumps({'code': 0, 'data': b64, 'file': f"师助AI_{body.get('title','')}.docx"})
                    print(f'Done: {len(docx_bytes)} bytes')
                elif doc_type == 'parse':
                    text = parse_file(body)
                    resp = json.dumps({'code': 0, 'data': text})
                elif doc_type == 'edit-excel':
                    b64_data, err = edit_excel(body)
                    if err:
                        resp = json.dumps({'code': -1, 'error': err})
                    else:
                        resp = json.dumps({'code': 0, 'data': b64_data, 'file': 'edited.xlsx'})
                elif doc_type == 'chat':
                    import urllib.request, urllib.error
                    # 统一使用阿里百炼：研究、故事板、文档文本与视觉能力都在同一工作空间内。
                    api_key = os.environ.get('AI_API_KEY', '').strip()
                    if not api_key:
                        resp = json.dumps({'code': -1, 'error': '服务器未配置阿里 AI_API_KEY，请在Railway环境变量中设置 AI_API_KEY'})
                    else:
                        messages = body.get('messages', [])
                        model = body.get('model', 'qwen-plus')
                        if str(model).startswith('deepseek'):
                            model = 'qwen-plus'
                        temp = body.get('temperature', 0.7)
                        max_tok = body.get('max_tokens', 4096)
                        post_data = json.dumps({'model': model, 'messages': messages, 'temperature': temp, 'max_tokens': max_tok, 'stream': False}).encode()
                        req = urllib.request.Request(
                            'https://ws-5ol6m5p8f4hikz1a.cn-beijing.maas.aliyuncs.com/compatible-mode/v1/chat/completions',
                            data=post_data,
                            headers={'Authorization': 'Bearer ' + api_key, 'Content-Type': 'application/json'}
                        )
                        try:
                            r = urllib.request.urlopen(req, timeout=120)
                            result = json.loads(r.read().decode())
                            reply = result['choices'][0]['message']['content']
                            resp = json.dumps({'code': 0, 'data': reply})
                            print(f'Chat reply: {len(reply)} chars')
                        except urllib.error.HTTPError as e:
                            resp = json.dumps({'code': -1, 'error': 'API错误: ' + e.read().decode()[:200]})
                        except Exception as e:
                            resp = json.dumps({'code': -1, 'error': str(e)[:200]})
                elif doc_type == 'search':
                    api_key = os.environ.get('AI_API_KEY', '').strip()
                    if not api_key:
                        resp = json.dumps({'code': -1, 'error': '联网搜索功能未配置（需在Railway设置 AI_API_KEY）'})
                    else:
                        import urllib.request, urllib.error
                        query = body.get('query', '')
                        search_data = json.dumps({
                            'model': 'qwen-plus',
                            'messages': [{'role': 'user', 'content': query}],
                            'enable_search': True,
                            'search_options': {'forced_search': True}
                        }).encode()
                        req = urllib.request.Request(
                            'https://ws-5ol6m5p8f4hikz1a.cn-beijing.maas.aliyuncs.com/compatible-mode/v1/chat/completions',
                            data=search_data,
                            headers={'Authorization': 'Bearer ' + api_key, 'Content-Type': 'application/json'}
                        )
                        try:
                            r = urllib.request.urlopen(req, timeout=120)
                            result = json.loads(r.read().decode())
                            reply = result['choices'][0]['message']['content']
                            resp = json.dumps({'code': 0, 'data': reply})
                        except Exception as e:
                            resp = json.dumps({'code': -1, 'error': '搜索失败: ' + str(e)[:200]})
                elif doc_type == 'ppt-build':
                    try:
                        deck = build_native_deck(body)
                        pptx_bytes = render_native_deck(deck)
                        result_label = 'native-pptxgen'
                        fid = str(uuid.uuid4())
                        fpath = os.path.join(DL_DIR, fid + '.pptx')
                        with open(fpath, 'wb') as fh:
                            fh.write(pptx_bytes)
                        if body.get('debug'):
                            print('[PPT TRACE] 10_ppt_result: ' + json.dumps({
                                'engine': result_label,
                                'slide_count': len(deck['slides']),
                                'pptx_bytes': len(pptx_bytes),
                            }, ensure_ascii=False))
                        resp = json.dumps({
                            'code': 0,
                            'url': '/dl/' + fid,
                            'title': deck['title'],
                            'slides': deck['slides'],
                            'debug': deck.get('debug') if body.get('debug') else None,
                        }, ensure_ascii=False)
                        print('[PPT API] ' + result_label + ' PPT exported:', deck['title'], len(deck['slides']))
                    except Exception as exc:
                        print('[PPT AGENT] failed:', repr(exc))
                        resp = json.dumps({'code': -1, 'error': str(exc)[:600]}, ensure_ascii=False)
                elif doc_type == 'ppt-research':
                    try:
                        pack = build_agent_research(body.get('topic', ''), body.get('detail', '中'))
                        resp = json.dumps({'code': 0, 'data': pack}, ensure_ascii=False)
                    except Exception as e:
                        resp = json.dumps({'code': -1, 'error': 'PPT联网研究失败: ' + str(e)[:300]}, ensure_ascii=False)
                elif doc_type == 'vision':
                    api_key = os.environ.get('AI_API_KEY', '').strip()
                    if not api_key:
                        resp = json.dumps({'code': -1, 'error': '图片识别功能未配置（需在Railway设置 AI_API_KEY）'})
                    else:
                        import urllib.request, urllib.error
                        image_b64 = body.get('image', '')
                        prompt = body.get('prompt', '请描述这张图片')
                        vision_data = json.dumps({
                            'model': 'qwen3-vl-flash',
                            'input': {
                                'messages': [{'role': 'user', 'content': [
                                    {'image': f'data:image/jpeg;base64,{image_b64}'},
                                    {'text': prompt}
                                ]}]
                            }
                        }).encode()
                        req = urllib.request.Request(
                            'https://dashscope.aliyuncs.com/api/v1/services/aigc/multimodal-generation/generation',
                            data=vision_data,
                            headers={'Authorization': 'Bearer ' + api_key, 'Content-Type': 'application/json'}
                        )
                        try:
                            r = urllib.request.urlopen(req, timeout=120)
                            result = json.loads(r.read().decode())
                            reply = result['output']['choices'][0]['message']['content']
                            resp = json.dumps({'code': 0, 'data': reply})
                        except Exception as e:
                            resp = json.dumps({'code': -1, 'error': '识别失败: ' + str(e)[:200]})
                elif doc_type == 'generate-image':
                    api_key = os.environ.get('AI_API_KEY', '').strip()
                    if not api_key:
                        resp = json.dumps({'code': -1, 'error': 'AI生图功能未配置（需在Railway设置 AI_API_KEY）'})
                    else:
                        import urllib.request, urllib.error, time
                        prompt = body.get('prompt', '')
                        img_data = json.dumps({
                            'model': 'z-image-turbo',
                            'input': {'prompt': prompt},
                            'parameters': {'size': '1024*1024', 'n': 1}
                        }).encode()
                        req = urllib.request.Request(
                            'https://ws-5ol6m5p8f4hikz1a.cn-beijing.maas.aliyuncs.com/api/v1/services/aigc/text2image/image-synthesis',
                            data=img_data,
                            headers={'Authorization': 'Bearer ' + api_key, 'Content-Type': 'application/json'}
                        )
                        try:
                            r = urllib.request.urlopen(req, timeout=120)
                            result = json.loads(r.read().decode())
                            task_id = result.get('output', {}).get('task_id', '')
                            if task_id:
                                # Poll for result
                                for _ in range(30):
                                    time.sleep(2)
                                    poll_req = urllib.request.Request(
                                        f'https://ws-5ol6m5p8f4hikz1a.cn-beijing.maas.aliyuncs.com/api/v1/tasks/{task_id}',
                                        headers={'Authorization': 'Bearer ' + api_key}
                                    )
                                    poll_r = urllib.request.urlopen(poll_req, timeout=30)
                                    poll_result = json.loads(poll_r.read().decode())
                                    status = poll_result.get('output', {}).get('task_status', '')
                                    if status == 'SUCCEEDED':
                                        img_url = poll_result.get('output', {}).get('results', [{}])[0].get('url', '')
                                        resp = json.dumps({'code': 0, 'data': img_url, 'format': 'url'})
                                        break
                                    elif status in ('FAILED', 'CANCELED'):
                                        resp = json.dumps({'code': -1, 'error': '生图失败'})
                                        break
                                else:
                                    resp = json.dumps({'code': -1, 'error': '生图超时'})
                            else:
                                resp = json.dumps({'code': -1, 'error': '提交生图任务失败'})
                        except Exception as e:
                            resp = json.dumps({'code': -1, 'error': '生图失败: ' + str(e)[:200]})
                elif doc_type == 'handwriting':
                    api_key = os.environ.get('AI_API_KEY', '').strip()
                    if not api_key:
                        resp = json.dumps({'code': -1, 'error': '\u9700\u8981\u914d\u7f6eAI_API_KEY'})
                    else:
                        import urllib.request, urllib.error, base64 as b64mod
                        image_b64 = body.get('image', '')
                        # Step 1: OCR using qwen-vl-max
                        ocr_data = json.dumps({
                            'model': 'qwen3-vl-flash',
                            'input': {
                                'messages': [{'role': 'user', 'content': [
                                    {'image': f'data:image/jpeg;base64,{image_b64}'},
                                    {'text': '\u8bf7\u8bc6\u522b\u8fd9\u5f20\u56fe\u7247\u4e2d\u7684\u6240\u6709\u6587\u5b57\u3002\u5982\u679c\u662f\u8868\u683c\u6570\u636e\uff0c\u7528\u5236\u8868\u7b26\t\u5206\u9694\u5217\uff0c\u6bcf\u884c\u4e00\u6761\u8bb0\u5f55\u3002\u5982\u679c\u662f\u6bb5\u843d\u6587\u5b57\uff0c\u76f4\u63a5\u8f93\u51fa\u6587\u5b57\u5185\u5bb9\u3002'}
                                ]}]       # close content[], msg_obj{}, msg_array[]
                            }              # close input{}
                        }).encode()
                        req = urllib.request.Request(
                            'https://dashscope.aliyuncs.com/api/v1/services/aigc/multimodal-generation/generation',
                            data=ocr_data,
                            headers={'Authorization': 'Bearer ' + api_key, 'Content-Type': 'application/json'}
                        )
                        try:
                            r = urllib.request.urlopen(req, timeout=120)
                            result = json.loads(r.read().decode())
                            text = result['output']['choices'][0]['message']['content']
                            # Step 2: Detect if table or text
                            detect_data = json.dumps({
                                'model': 'qwen-turbo',
                                'messages': [{'role': 'user', 'content': f'\u5224\u65ad\u4ee5\u4e0b\u5185\u5bb9\u662f\u8868\u683c\u6570\u636e\u8fd8\u662f\u6bb5\u843d\u6587\u5b57\uff0c\u53ea\u56de\u7b54"table"\u6216"text"\uff1a\n{text[:1500]}'}],
                                'temperature': 0.1
                            }).encode()
                            detect_req = urllib.request.Request(
                                'https://ws-5ol6m5p8f4hikz1a.cn-beijing.maas.aliyuncs.com/compatible-mode/v1/chat/completions',
                                data=detect_data,
                                headers={'Authorization': 'Bearer ' + api_key, 'Content-Type': 'application/json'}
                            )
                            detect_r = urllib.request.urlopen(detect_req, timeout=30)
                            detect_result = json.loads(detect_r.read().decode())
                            doc_type_detected = detect_result['choices'][0]['message']['content'].strip().lower()
                            
                            if 'table' in doc_type_detected:
                                # Generate Excel
                                import openpyxl
                                wb = openpyxl.Workbook()
                                ws = wb.active
                                ws.title = '\u6570\u636e'
                                lines = text.strip().split('\n')
                                for ri, line in enumerate(lines):
                                    cells = line.split('\t')
                                    for ci, cell in enumerate(cells):
                                        ws.cell(ri+1, ci+1, cell.strip())
                                buf = io.BytesIO()
                                wb.save(buf)
                                fid = str(uuid.uuid4())
                                fpath = os.path.join(DL_DIR, fid + '.xlsx')
                                with open(fpath, 'wb') as fh: fh.write(buf.getvalue())
                                resp = json.dumps({'code': 0, 'url': '/dl/' + fid, 'detected': 'table', 'ocrText': text[:5000]})
                                print(f'Handwriting: detected as table, {len(lines)} rows')
                            else:
                                # Generate Word
                                w_data = {'title': '\u624b\u5199\u6587\u6863', 'content': text}
                                docx_bytes = gen_docx(w_data)
                                fid = str(uuid.uuid4())
                                fpath = os.path.join(DL_DIR, fid + '.docx')
                                with open(fpath, 'wb') as fh: fh.write(docx_bytes)
                                resp = json.dumps({'code': 0, 'url': '/dl/' + fid, 'detected': 'text', 'ocrText': text[:5000]})
                                print(f'Handwriting: detected as text, {len(text)} chars')
                        except urllib.error.HTTPError as e:
                            body = e.read().decode()[:300]
                            resp = json.dumps({'code': -1, 'error': '\u8bc6\u522b\u5931\u8d25 (' + str(e.code) + '): ' + body})
                        except Exception as e:
                            resp = json.dumps({'code': -1, 'error': str(e)[:200]})
                else:
                    print(f'Generating PPT: {body.get("title","")}')
                    pptx_bytes = gen(body)
                    fid = str(uuid.uuid4())
                    fpath = os.path.join(DL_DIR, fid + '.pptx')
                    with open(fpath, 'wb') as fh:
                        fh.write(pptx_bytes)
                    resp = json.dumps({'code': 0, 'url': '/dl/' + fid})
                    print(f'Done: {len(pptx_bytes)} bytes, url=/dl/{fid}')
                
                self.send_response(200)
                self.send_header('Content-Type', 'application/json')
                self._set_cors()
                self.end_headers()
                self.wfile.write(resp.encode())
            except Exception as e:
                # 线上 500 必须留下完整栈，才能区分研究、讲述稿、大纲、图片或渲染哪个阶段失败。
                import traceback
                print('[PPT API] request failed:', repr(e))
                traceback.print_exc()
                self.send_response(500)
                self.send_header('Content-Type', 'application/json')
                self._set_cors()
                self.end_headers()
                self.wfile.write(json.dumps({'code': -1, 'error': str(e)}).encode())

        def log_message(self, fmt, *args):
            print(f'[PPT API] {args[0] if args else fmt}')
    return Handler

if __name__ == '__main__':
    PORT = int(os.environ.get('PORT', 8080))
    server = http.server.HTTPServer(('0.0.0.0', PORT), make_handler())
    print(f'师助AI PPT API running on port {PORT}')
    server.serve_forever()
